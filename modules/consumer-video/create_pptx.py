#!/usr/bin/env python3
"""
Create PowerPoint presentation directly using python-pptx
"""

import json
import subprocess
import sys
from pathlib import Path
from datetime import datetime

from config import PATHS

# Auto-install python-pptx if needed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

# Load analysis
PROCESSED_DIR = Path(PATHS["processed"])
ANALYSIS_PATH = PROCESSED_DIR / "jtbd_analysis.json"

def load_analysis():
    with open(ANALYSIS_PATH, 'r') as f:
        return json.load(f)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullets, notes=""):
    """Add a content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    # Add bullets
    content = slide.placeholders[1].text_frame
    content.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            content.text = bullet
        else:
            p = content.add_paragraph()
            p.text = bullet
            p.level = 0

    # Add speaker notes if provided
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def add_table_slide(prs, title, headers, rows, notes=""):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.75)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True

    # Add table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5 * rows_count)

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    # Set column width
    for col_idx in range(cols_count):
        table.columns[col_idx].width = Inches(9 / cols_count)

    # Add headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # Add rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Add speaker notes if provided
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide

def create_presentation(analysis):
    """Create the full PowerPoint presentation"""
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "3M Lighting Consumer Research",
        "Jobs-to-be-Done Analysis\nRungs of the Ladder Framework"
    )

    # Slide 2: Research Overview
    add_content_slide(
        prs,
        "Research Overview",
        [
            f"{analysis['summary']['total_consumers']} consumers analyzed across 82 video interviews",
            f"{analysis['summary']['total_jobs']} distinct jobs extracted from verbatim quotes",
            "Framework: Jobs-to-be-Done 'Rungs of the Ladder'",
            "Focus: Understanding not just what consumers do, but why they care"
        ],
        notes="Emphasize depth of analysis - real consumer voices, not survey data. All insights are directly quoted and verifiable."
    )

    # Slide 3: Rungs of the Ladder Framework
    add_content_slide(
        prs,
        "The Rungs of the Ladder Approach",
        [
            "Bottom Rung - Functional Jobs: What consumers are trying to accomplish",
            "  Example: Install accent lighting, illuminate artwork",
            "",
            "Middle Rungs - Emotional Jobs: How consumers want to feel",
            "  Example: Feel capable, keep it simple, avoid regret",
            "",
            "Top Rungs - Social Jobs: How consumers want to be perceived",
            "  Example: Appear modern/sophisticated, demonstrate competence"
        ],
        notes="Higher rungs reveal the true value drivers. Consumers don't just want lighting - they want to feel capable and appear resourceful."
    )

    # Slide 4: Functional Jobs
    add_content_slide(
        prs,
        "Functional Jobs - What They're Trying to Do",
        [
            "Control Lighting (7 consumers): Remote, timers, motion sensors",
            "Avoid Cords/Wiring (9 consumers): Battery, rechargeable solutions",
            "Illuminate Space (5 consumers): Light up art, features",
            "Highlight Features (3 consumers): Draw attention to design elements",
            "",
            "Consumer Voice:",
            "\"I didn't want wires hanging on walls... no ability to create electrical work inside walls\" - GeneK"
        ],
        notes="Wireless/battery solutions are table stakes, not differentiators. Control features serve multiple jobs simultaneously."
    )

    # Slide 5: Emotional Jobs
    add_content_slide(
        prs,
        "Emotional Jobs - How They Want to Feel",
        [
            "Keep It Simple (7 consumers): Avoid hassle, maintain control",
            "Achieve Desired Aesthetic (6 consumers): Modern, luxury, elegant look",
            "Avoid High Costs (4 consumers): \"$1000 every time electrician comes\"",
            "Feel Capable (6 consumers): Pride in DIY accomplishment",
            "Avoid Regret (3 consumers): Get it right the first time",
            "",
            "Consumer Voice:",
            "\"It still gives me that modern luxury aspect... didn't have to invest much money\" - TylrD"
        ],
        notes="Emotional jobs often drive more value than functional jobs. Aesthetic is non-negotiable, but must be achievable."
    )

    # Slide 6: Social Jobs
    add_content_slide(
        prs,
        "Social Jobs - How They Want to Be Perceived",
        [
            "Be Seen as Creative (7 consumers): Design-conscious choices",
            "Showcase Home (5 consumers): Impress visitors, display identity",
            "Appear Modern/Sophisticated (5 consumers): Tech-savvy, upscale",
            "Display Personal Identity (3 consumers): Family photos, art collection",
            "Demonstrate DIY Competence: \"No need to hire\"",
            "",
            "Consumer Voice:",
            "\"Motion detected... still get that same look as modern and classic design\" - TylrD"
        ],
        notes="Social signaling matters - lighting is a design statement. DIY capability creates pride when achievable but not trivial."
    )

    # Slide 7: Opportunity Map
    add_content_slide(
        prs,
        "Opportunity Map: Ease of Implementation vs Aspirational Value",
        [
            "HIGH Ease, HIGH Aspirational (Sweet Spot):",
            "  • Wireless Premium Positioning",
            "  • Smart Control as Standard",
            "  • Design-Forward DIY",
            "",
            "Gap in Market: Current options force choice between:",
            "  • Easy install (but basic/ugly) OR",
            "  • High-end look (but expensive/difficult)",
            "",
            "3M Opportunity: Premium DIY without sacrificing aesthetics"
        ],
        notes="Sweet spot is HIGH Ease + HIGH Aspirational quadrant. Current market forces consumers to choose between easy install and high-end look."
    )

    # Slide 8: Jobs Summary Table
    add_table_slide(
        prs,
        "Jobs Summary by Priority",
        ["Job", "Type", "Consumers", "Priority"],
        [
            ["Avoid Cords/Wiring", "Functional", "9", "HIGH"],
            ["Control Lighting", "Functional", "7", "HIGH"],
            ["Keep It Simple", "Emotional", "7", "HIGH"],
            ["Be Seen as Creative", "Social", "7", "HIGH"],
            ["Achieve Desired Aesthetic", "Emotional", "6", "HIGH"],
            ["Feel Capable", "Emotional", "6", "HIGH"],
            ["Showcase Home", "Social", "5", "MEDIUM"],
            ["Appear Modern/Sophisticated", "Social", "5", "MEDIUM"],
            ["Illuminate Space", "Functional", "5", "MEDIUM"]
        ],
        notes="Priority based on consumer diversity and mention frequency. Jobs with 5+ consumers represent validated opportunities."
    )

    # Slide 9: Recommended Focus Areas
    add_content_slide(
        prs,
        "Recommended Focus Areas",
        [
            "1. Wireless Premium Positioning:",
            "   Premium battery solutions without sacrificing aesthetic",
            "   Serves: Avoid electrician + Look finished + Appear modern",
            "",
            "2. Smart Control as Standard:",
            "   Make remote, timers, dimming standard features",
            "   Serves: Control flexibility + Feel capable + Signal sophistication",
            "",
            "3. Installation Confidence System:",
            "   Tools/guides for right product selection and placement",
            "   Serves: Feel capable + Avoid regret + Demonstrate competence",
            "",
            "4. Design-Forward DIY:",
            "   Position as choice for design-conscious DIYers",
            "   Serves: Achieve aesthetic + Be creative + Showcase home"
        ],
        notes="Each opportunity serves multiple ladder rungs simultaneously. Gap exists between contractor-grade ugly and designer-grade expensive."
    )

    # Slide 10: Key Insights
    add_content_slide(
        prs,
        "Key Insights & Patterns",
        [
            "Pattern 1: The Electrician Avoidance",
            "  9 of 15 consumers explicitly avoided electricians",
            "  Not just cost - it's about control, speed, DIY pride",
            "",
            "Pattern 2: Control as Feature Multiplier",
            "  Remote, timers, motion sensors serve multiple jobs",
            "  Functional + Emotional + Social value simultaneously",
            "",
            "Pattern 3: Regret is Real",
            "  3 consumers expressed regret about choices",
            "  Opportunity for better upfront guidance",
            "",
            "Pattern 4: Aesthetic Non-Negotiable",
            "  6 consumers used 'modern,' 'luxury,' 'elegant'",
            "  Lighting is a design statement, not just illumination"
        ],
        notes="These patterns reveal unmet needs and innovation opportunities."
    )

    # Slide 11: Next Steps
    add_content_slide(
        prs,
        "Next Steps",
        [
            "Immediate (2 weeks):",
            "  • Internal workshop with product and marketing teams",
            "  • Opportunity prioritization vote",
            "  • Competitive mapping against opportunity matrix",
            "",
            "Short Term (4-6 weeks):",
            "  • Concept development for high-priority jobs",
            "  • Consumer validation with 20-30 participants",
            "  • Messaging framework addressing emotional/social jobs",
            "",
            "Medium Term (3 months):",
            "  • Prototype development for top 2-3 opportunities",
            "  • In-home testing and extended trials",
            "  • Go-to-market planning based on ladder insights"
        ],
        notes="All 82 videos and transcripts available for deeper exploration. Source files verifiable for any insight cited."
    )

    return prs

def main():
    print("Loading JTBD analysis...")
    analysis = load_analysis()

    print("Creating PowerPoint presentation...")
    prs = create_presentation(analysis)

    # Save with timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_dir = Path(PATHS["outputs"])
    output_dir.mkdir(parents=True, exist_ok=True)
    filename = output_dir / f"3M_Lighting_JTBD_Presentation_{timestamp}.pptx"
    prs.save(filename)

    print(f"\n✓ PowerPoint created successfully!")
    print(f"  File: {filename}")
    print(f"  Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
