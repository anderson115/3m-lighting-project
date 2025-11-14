#!/usr/bin/env python3
"""
Use Offbrain template MASTER LAYOUTS (not custom slides) to create JTBD presentation
"""

from config import PATHS

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE_PATH = Path("/Users/anderson115/Desktop/Offbrain_FINAL.pptx")
OUTPUT_PATH = Path(PATHS["outputs"]) / "3M_Lighting_JTBD_Final_Complete-Codex.pptx"

def create_presentation():
    """Create presentation using Offbrain MASTER layouts"""

    # Load template
    prs = Presentation(str(TEMPLATE_PATH))

    # Delete ALL existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Slide 1: Title Slide (Layout 0)
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "3M Lighting"
    slide.placeholders[1].text = "Jobs-to-be-Done Analysis\nConsumer Insights & Strategic Opportunities"

    # Use Layout 2 (Section Header) - has BODY type placeholder for bullets
    content_layout = prs.slide_layouts[2]  # "Section Header" with BODY placeholder

    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Executive Summary"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "15 consumers analyzed across 87 video transcripts"

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "Consumers are creating gallery-worthy lighting without professional crews:"

        p = tf.add_paragraph()
        p.text = "Rechargeable spotlights & LED strips showcase art and family stories"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Adhesive mounts keep floors clear while maintaining designer aesthetics"
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "Three High-Value Opportunity Clusters:"
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = "1. Precision DIY gallery systems with alignment tools"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "2. Adaptive ambiance controls (retrofit-friendly)"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "3. Low-profile safety-plus-style lighting for pathways"
        p.level = 1

    # Slide 3: Methodology
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Methodology: Jobs-to-be-Done Framework"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "Rungs of the Ladder Approach"
        tf.paragraphs[0].font.bold = True

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "FUNCTIONAL JOBS (Bottom): Outcomes to achieve"

        p = tf.add_paragraph()
        p.text = "Example: Illuminate focal art without hiring electricians"
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "EMOTIONAL JOBS (Middle): Feelings to experience"

        p = tf.add_paragraph()
        p.text = "Example: Feel proud delivering designer-level results"
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "SOCIAL JOBS (Top): Perceptions to cultivate"

        p = tf.add_paragraph()
        p.text = "Example: Impress guests with gallery-like presentation"
        p.level = 1

    # Slide 4: Functional Jobs
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Functional Jobs (Bottom Rung)"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "7 Core Functional Jobs Identified"
        tf.paragraphs[0].font.bold = True

        jobs = [
            "F1. Illuminate focal art and architectural features",
            "F2. Add accent lighting avoiding electricians/wiring",
            "F3. Control ambiance with adjustable modes",
            "F4. Keep floors and walls uncluttered",
            "F5. Build expandable gallery lighting",
            "F6. Deliver customizable color experiences",
            "F7. Provide multipurpose accent + safety lighting"
        ]

        for job in jobs:
            p = tf.add_paragraph()
            p.text = job

    # Slide 5: F1 & F2 with quotes
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Functional Jobs: Art Illumination & DIY Install"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "F1: Illuminate focal art so it commands attention"
        tf.paragraphs[0].font.bold = True

        p = tf.add_paragraph()
        p.text = '"We wanted accent lighting to draw attention to the fireplace mantel, as well as the new art that hangs right above it." — GeneK'
        p.level = 1
        p.font.italic = True

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "F2: Avoid hiring professionals or running new wiring"
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = '"I didn\'t want any wires hanging on walls... I don\'t have any ability to create electrical work inside the walls." — GeneK'
        p.level = 1
        p.font.italic = True

    # Slide 6: Emotional Jobs
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Emotional Jobs (Middle Rungs)"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "5 Core Emotional Jobs"
        tf.paragraphs[0].font.bold = True

        ejobs = [
            "E1. Feel proud and accomplished delivering designer results",
            "E2. Feel calm, relaxed, and connected to the space",
            "E3. Feel confident navigating DIY installs",
            "E4. Reduce stress about overspending",
            "E5. Feel ready to celebrate evolving family stories"
        ]

        for job in ejobs:
            p = tf.add_paragraph()
            p.text = job

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = '"I felt like a proud dad to do this for my son." — WilliamS'
        p.level = 1
        p.font.italic = True

    # Slide 7: Social Jobs
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Social Jobs (Top Rungs)"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "4 Core Social Jobs"
        tf.paragraphs[0].font.bold = True

        sjobs = [
            "S1. Impress guests with gallery-like presentation",
            "S2. Showcase personal and family narratives",
            "S3. Demonstrate resourcefulness and DIY savvy",
            "S4. Keep home looking polished and magazine-ready"
        ]

        for job in sjobs:
            p = tf.add_paragraph()
            p.text = job

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = '"When I have guests over, I like to highlight my artwork in my house with the spotlight." — FarahN'
        p.level = 1
        p.font.italic = True

    # Slide 8: Opportunity Map
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Opportunity Map: Priority Zones"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "Matrix: Implementation Effort vs. Aspirational Lift"
        tf.paragraphs[0].font.bold = True

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "HIGH PRIORITY (Low-Moderate Effort, High Aspirational):"
        p.font.bold = True

        priorities = [
            "F5: Expandable gallery lighting (modular systems)",
            "F6: Customizable color experiences",
            "E2: Calm and mood-setting solutions",
            "All social jobs (gallery presentation, storytelling)"
        ]

        for item in priorities:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1

    # Slide 9: Key Insights
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Key Insights & Patterns"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame

        insights = [
            ("1. Adhesive resilience is a weak link", "Outdoor heat and textured surfaces challenge current tapes"),
            ("2. Precision layout tools are missing", "Consumers Googling spacing formulas, want even-spacing tools"),
            ("3. Consumers optimize infrastructure first", "Tap attic circuits, cut holes, pivot to battery before calling pros"),
            ("4. Control drives emotional payoff", "Remote temperature, dimmers, motion = 'feels finished'"),
            ("5. Gallery storytelling is differentiator", "Families build evolving walls with backup lights")
        ]

        tf.text = insights[0][0]
        tf.paragraphs[0].font.bold = True
        p = tf.add_paragraph()
        p.text = insights[0][1]
        p.level = 1

        for title, detail in insights[1:]:
            p = tf.add_paragraph()
            p.text = ""
            p = tf.add_paragraph()
            p.text = title
            p.font.bold = True
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1

    # Slide 10: Recommendations
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Recommended Focus Areas"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame

        tf.text = "1. Gallery Alignment & Adhesive Assurance Kit"
        tf.paragraphs[0].font.bold = True

        p = tf.add_paragraph()
        p.text = "Heat-rated adhesive + laser/spacing guides"
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "2. Rechargeable Ambiance Control Bundle"
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = "Tunable lights + dimmer remotes + presets"
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "3. Hybrid Décor-Safety Bar Lighting"
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = "Switchable gallery/path modes"
        p.level = 1

    # Slide 11: Next Steps
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Next Steps"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "IMMEDIATE (2-4 weeks):"
        tf.paragraphs[0].font.bold = True

        p = tf.add_paragraph()
        p.text = "Workshop with product and marketing teams"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Prioritization vote on opportunity clusters"
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "SHORT TERM (4-8 weeks):"
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = "Develop alignment kit prototypes"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Test ambiance control bundles with consumers"
        p.level = 1

        p = tf.add_paragraph()
        p.text = ""

        p = tf.add_paragraph()
        p.text = "MEDIUM TERM (3 months):"
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = "Launch gallery system pilot"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Build marketing narratives around emotional/social jobs"
        p.level = 1

    # Save
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH

def main():
    print("Creating presentation using Offbrain MASTER layouts...")
    output = create_presentation()
    print(f"\n✓ COMPLETE!")
    print(f"  File: {output}")
    print(f"  Using Layout 2 (Section Header) with BODY placeholder")

if __name__ == "__main__":
    main()
