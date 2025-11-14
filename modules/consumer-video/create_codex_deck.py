#!/usr/bin/env python3
"""
Create polished client-ready JTBD presentation using Offbrain template
Based on comprehensive analysis with proper citations
"""

from config import PATHS

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
TEMPLATE_PATH = Path("/Users/anderson115/Desktop/Offbrain_FINAL.pptx")
OUTPUT_PATH = Path(PATHS["outputs"]) / "3M_Lighting_JTBD_Final_Complete-Codex.pptx"

def add_slide_with_title(prs, layout_idx, title_text):
    """Add a slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if hasattr(slide.shapes, 'title') and slide.shapes.title:
        slide.shapes.title.text = title_text
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False):
    """Add a text box to slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.name = 'Calibri'
        paragraph.font.bold = bold
    return txBox

def add_bullets(slide, items, left=None, top=None):
    """Add bullet points to slide"""
    if left is None:
        left = Inches(0.75)
    if top is None:
        top = Inches(1.75)

    width = Inches(8.5)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.font.size = Pt(11) if item.strip() and not item.startswith('**') else Pt(9)
        p.font.name = 'Calibri'

        # Handle bold markers
        if item.startswith('**') and item.endswith('**'):
            p.text = item.strip('**')
            p.font.bold = True
            p.font.size = Pt(13)

        # Set indentation levels
        if item.startswith('    '):
            p.level = 2
        elif item.startswith('  '):
            p.level = 1
        else:
            p.level = 0

def create_presentation():
    """Create the full presentation"""

    # Load template
    prs = Presentation(str(TEMPLATE_PATH))

    # Use blank layout for all content slides
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if hasattr(slide.shapes, 'title'):
        slide.shapes.title.text = "3M Lighting"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Jobs-to-be-Done Analysis\nConsumer Insights & Strategic Opportunities"

    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Executive Summary", font_size=32, bold=True)

    add_bullets(slide, [
        "**15 consumers analyzed across 87 video transcripts**",
        "",
        "Consumers are creating gallery-worthy lighting without professional crews:",
        "• Rechargeable spotlights & LED strips showcase art and family stories",
        "• Adhesive mounts keep floors clear while maintaining designer aesthetics",
        "• Remote controls, dimmers, and color programs deliver emotional payoffs",
        "",
        "**Three High-Value Opportunity Clusters:**",
        "1. Precision DIY gallery systems with alignment tools",
        "2. Adaptive ambiance controls (retrofit-friendly)",
        "3. Low-profile safety-plus-style lighting for pathways"
    ], left=Inches(0.75), top=Inches(1.5))

    # Slide 3: Methodology
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Methodology: Jobs-to-be-Done Framework", font_size=28, bold=True)

    add_bullets(slide, [
        "**Rungs of the Ladder Approach:**",
        "",
        "• FUNCTIONAL JOBS (Bottom): Outcomes to achieve",
        "  Example: Illuminate focal art without hiring electricians",
        "",
        "• EMOTIONAL JOBS (Middle): Feelings to experience",
        "  Example: Feel proud delivering designer-level results",
        "",
        "• SOCIAL JOBS (Top): Perceptions to cultivate",
        "  Example: Impress guests with gallery-like presentation",
        "",
        "**Data Sources:**",
        "• 87 transcript files covering 15 unique consumers",
        "• 30+ interview modules (Activities 1-10 + walk-throughs)",
        "• All jobs validated with verbatim quotes and citations"
    ])

    # Slide 4: Functional Jobs Overview
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Functional Jobs (Bottom Rung)", font_size=28, bold=True)

    add_bullets(slide, [
        "**7 Core Functional Jobs Identified:**",
        "",
        "F1. Illuminate focal art and architectural features",
        "F2. Add accent lighting avoiding electricians/wiring",
        "F3. Control ambiance with adjustable modes",
        "F4. Keep floors and walls uncluttered",
        "F5. Build expandable gallery lighting",
        "F6. Deliver customizable color experiences",
        "F7. Provide multipurpose accent + safety lighting",
        "",
        "**Key Pattern:**",
        "Consumers prioritize DIY solutions that avoid contractor costs",
        "while achieving professional-quality results"
    ])

    # Slide 5: F1 & F2 Detail
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Functional Jobs: Art Illumination & DIY Install", font_size=24, bold=True)

    add_bullets(slide, [
        "**F1: Illuminate focal art so it commands attention**",
        "",
        "\"We wanted accent lighting to draw attention to the fireplace",
        "mantel, as well as the new art that hangs right above it.\"",
        "— GeneK",
        "",
        "**F2: Avoid hiring professionals or running new wiring**",
        "",
        "\"I didn't want any wires hanging on walls... I don't have any",
        "ability to create electrical work inside the walls.\"",
        "— GeneK",
        "",
        "\"Every time you get somebody to come into your house,",
        "it's a thousand dollars... No need to hire expensive installer.\"",
        "— FarahN"
    ])

    # Slide 6: F3 & F4 Detail
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Functional Jobs: Control & Clean Aesthetics", font_size=24, bold=True)

    add_bullets(slide, [
        "**F3: Control ambiance with adjustable lighting modes**",
        "",
        "\"I can change the temperature and they will highlight this painting.\"",
        "— GeneK (remote-controlled accent lights)",
        "",
        "\"All my lights are on dimmers.\"",
        "— MarkR",
        "",
        "**F4: Keep floors and walls uncluttered**",
        "",
        "\"It lights up the mural, and that way it leaves the ground",
        "clear so nobody trips over it.\"",
        "— CarrieS",
        "",
        "\"They are also motion activated, so they light up my hall.\"",
        "— RobinL"
    ])

    # Slide 7: Emotional Jobs
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Emotional Jobs (Middle Rungs)", font_size=28, bold=True)

    add_bullets(slide, [
        "**5 Core Emotional Jobs:**",
        "",
        "E1. Feel proud and accomplished delivering designer results",
        "E2. Feel calm, relaxed, and connected to the space",
        "E3. Feel confident navigating DIY installs",
        "E4. Reduce stress about overspending",
        "E5. Feel ready to celebrate evolving family stories",
        "",
        "**Key Insight:**",
        "\"I felt like a proud dad to do this for my son.\"",
        "— WilliamS",
        "",
        "\"Doing it myself... that feeling of completion.\"",
        "— ChristianL"
    ])

    # Slide 8: Social Jobs
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Social Jobs (Top Rungs)", font_size=28, bold=True)

    add_bullets(slide, [
        "**4 Core Social Jobs:**",
        "",
        "S1. Impress guests with gallery-like presentation",
        "S2. Showcase personal and family narratives",
        "S3. Demonstrate resourcefulness and DIY savvy",
        "S4. Keep home looking polished and magazine-ready",
        "",
        "**Consumer Voice:**",
        "",
        "\"When I have guests over, I like to highlight my artwork",
        "in my house with the spotlight.\"",
        "— FarahN",
        "",
        "\"Lighting is very important. It speaks volumes of your home.\"",
        "— MarkR"
    ])

    # Slide 9: Opportunity Map
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Opportunity Map: Priority Zones", font_size=28, bold=True)

    add_bullets(slide, [
        "**Matrix: Implementation Effort vs. Aspirational Lift**",
        "",
        "**HIGH PRIORITY (Low-Moderate Effort, High Aspirational):**",
        "• F5: Expandable gallery lighting (modular systems)",
        "• F6: Customizable color experiences",
        "• F7: Multipurpose accent + safety lighting",
        "• E2: Calm and mood-setting solutions",
        "• E5: Memory-ready lighting packages",
        "• S1-S4: All social jobs (gallery presentation, storytelling)",
        "",
        "**CONFIDENCE GAP (High Effort, High Aspirational):**",
        "• E3: DIY confidence - needs measurement tools,",
        "  electrical guidance, and templates",
        "",
        "**FOUNDATIONAL (Low Effort, Medium Aspirational):**",
        "• F2, F4, E4: Adhesives, cable mgmt, value-tier SKUs"
    ])

    # Slide 10: Key Insights
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Key Insights & Patterns", font_size=28, bold=True)

    add_bullets(slide, [
        "**1. Adhesive resilience is a weak link**",
        "Outdoor heat and textured surfaces challenge current tapes",
        "— CarrieS: fixtures fell in Arizona heat",
        "",
        "**2. Precision layout tools are missing**",
        "\"Kind of a mess\" - Googling spacing formulas",
        "— ChristianL: wants even-spacing tool",
        "",
        "**3. Consumers optimize existing infrastructure first**",
        "Tap attic circuits, cut pre-approved holes, pivot to battery",
        "— FarahN, AlysonT, AlanG",
        "",
        "**4. Control drives emotional payoff**",
        "Remote temperature, dimmers, motion sensors = \"feels finished\"",
        "",
        "**5. Gallery storytelling is a differentiator**",
        "Families build evolving walls with backup lights and frames"
    ])

    # Slide 11: Recommendations
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Recommended Focus Areas", font_size=28, bold=True)

    add_bullets(slide, [
        "**1. Gallery Alignment & Adhesive Assurance Kit**",
        "Heat-rated adhesive strips + laser/spacing guides + mounting rails",
        "Addresses: CarrieS adhesion failures, ChristianL spacing struggles",
        "",
        "**2. Rechargeable Ambiance Control Bundle**",
        "Tunable picture lights + dimmer remotes + preset warmth modes",
        "Addresses: MarkR, GeneK, TiffanyO control needs",
        "",
        "**3. Hybrid Décor-Safety Bar Lighting**",
        "Slim, magnet-based with switchable gallery/path modes",
        "Addresses: RobinL and CarrieS hallway/outdoor needs",
        "",
        "**4. DIY Confidence Companion**",
        "Guided install content + templates + smart tool recommendations",
        "Addresses: AlysonT, DianaL, AlanG confidence gaps",
        "",
        "**5. Modular Colour-Play LED Ecosystem**",
        "Design-forward kits with curated scenes + motion add-ons",
        "Addresses: TylrD, WilliamS, TiffanyO expressive goals"
    ])

    # Slide 12: Next Steps
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.75),
                 "Next Steps", font_size=32, bold=True)

    add_bullets(slide, [
        "**IMMEDIATE (2-4 weeks):**",
        "• Workshop with product and marketing teams",
        "• Prioritization vote on opportunity clusters",
        "• Prototype specifications for top 3 solutions",
        "",
        "**SHORT TERM (4-8 weeks):**",
        "• Develop alignment kit prototypes",
        "• Test ambiance control bundles with 10-15 consumers",
        "• Validate adhesive improvements for heat/texture",
        "",
        "**MEDIUM TERM (3 months):**",
        "• Launch gallery system pilot",
        "• In-home testing with consumer advocates",
        "• Build marketing narratives around emotional/social jobs",
        "",
        "**All insights traceable to 87 source transcripts**",
        "**Complete citation documentation available**"
    ])

    # Save
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH

def main():
    print("Creating polished Codex presentation...")
    print(f"Using template: {TEMPLATE_PATH}")

    output = create_presentation()

    print(f"\n✓ COMPLETE!")
    print(f"  File: {output}")
    print(f"  Slides: 12 (Title + 11 content)")
    print(f"\nSlide Structure:")
    print(f"  1. Title")
    print(f"  2. Executive Summary")
    print(f"  3. Methodology")
    print(f"  4. Functional Jobs Overview")
    print(f"  5. F1 & F2 Detail (Art + DIY)")
    print(f"  6. F3 & F4 Detail (Control + Clean)")
    print(f"  7. Emotional Jobs")
    print(f"  8. Social Jobs")
    print(f"  9. Opportunity Map")
    print(f"  10. Key Insights")
    print(f"  11. Recommendations")
    print(f"  12. Next Steps")

if __name__ == "__main__":
    main()
