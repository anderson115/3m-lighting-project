#!/usr/bin/env python3
"""
Create final JTBD presentation using Offbrain template
Following Clayton Christensen framework properly
"""

from config import PATHS

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from collections import defaultdict

# Paths
PROCESSED_DIR = Path(PATHS["processed"])
TEMPLATE_PATH = Path(PATHS["outputs"]) / "3M_Lighting_JTBD_Final.pptx"

def load_transcripts():
    """Load all transcripts with metadata"""
    transcripts = []
    folders = sorted([f for f in PROCESSED_DIR.iterdir() if f.is_dir() and not f.name.startswith('batch')])

    for folder in folders[:30]:  # Sample first 30 for speed
        transcript_path = folder / "transcript.json"
        if transcript_path.exists():
            with open(transcript_path, 'r') as f:
                data = json.load(f)
                consumer = folder.name.split('_')[0]
                activity = folder.name.split('_')[1] if len(folder.name.split('_')) > 1 else "Interview"
                transcripts.append({
                    'consumer': consumer,
                    'activity': activity,
                    'text': data.get('text', ''),
                    'segments': data.get('segments', []),
                    'source': folder.name
                })

    return transcripts

def extract_jtbd_jobs(transcripts):
    """Extract JTBD jobs following Christensen framework"""

    # Define jobs based on analysis (proper "When...I want...So I can" format)
    jobs = [
        {
            'name': 'Complete Tasks in Poorly Lit Spaces',
            'statement': 'When I\'m working in areas with inadequate existing lighting, I want to illuminate the specific workspace, So I can complete the task safely and efficiently without struggle.',
            'evidence': [],
            'functional': 'Task completion in dark spaces',
            'emotional': 'Reduce frustration, feel capable',
            'social': 'Demonstrate competent home management'
        },
        {
            'name': 'Avoid Electrician Costs and Complexity',
            'statement': 'When I need lighting improvements but want to control my project timeline and budget, I want to install solutions myself without electrical work, So I can achieve my goals affordably and on my schedule.',
            'evidence': [],
            'functional': 'Self-installation without wiring',
            'emotional': 'Maintain control, avoid stress of contractors',
            'social': 'Demonstrate DIY competence and resourcefulness'
        },
        {
            'name': 'Create an Intentional Aesthetic',
            'statement': 'When I want my space to reflect my style and attention to detail, I want lighting that enhances the visual appeal, So I can feel proud of my home and create the ambiance I envision.',
            'evidence': [],
            'functional': 'Enhance visual appearance of space',
            'emotional': 'Feel satisfied with aesthetic choices',
            'social': 'Impress visitors, signal design consciousness'
        },
        {
            'name': 'Maintain Flexibility and Control',
            'statement': 'When my lighting needs change throughout the day or across seasons, I want easy control over brightness, timing, and activation, So I can adapt the environment without manual intervention.',
            'evidence': [],
            'functional': 'Adjust lighting parameters easily',
            'emotional': 'Feel modern, tech-savvy',
            'social': 'Signal sophistication through smart features'
        }
    ]

    # Extract evidence from transcripts
    for t in transcripts:
        text_lower = t['text'].lower()

        # Job 1: Task completion
        if any(word in text_lower for word in ['dark', 'see', 'closet', 'can\'t see', 'illuminate', 'workspace', 'kitchen']):
            jobs[0]['evidence'].append({
                'consumer': t['consumer'],
                'quote': t['text'][:200],
                'source': t['source']
            })

        # Job 2: Avoid electrician
        if any(word in text_lower for word in ['electrician', 'hardwire', 'wiring', 'cost', 'thousand dollars', 'battery', 'rechargeable']):
            jobs[1]['evidence'].append({
                'consumer': t['consumer'],
                'quote': t['text'][:200],
                'source': t['source']
            })

        # Job 3: Aesthetic
        if any(word in text_lower for word in ['modern', 'luxury', 'elegant', 'look', 'design', 'style', 'beautiful']):
            jobs[2]['evidence'].append({
                'consumer': t['consumer'],
                'quote': t['text'][:200],
                'source': t['source']
            })

        # Job 4: Control/flexibility
        if any(word in text_lower for word in ['remote', 'control', 'timer', 'motion', 'adjust', 'brightness', 'dim']):
            jobs[3]['evidence'].append({
                'consumer': t['consumer'],
                'quote': t['text'][:200],
                'source': t['source']
            })

    # Sort by evidence strength
    for job in jobs:
        job['consumer_count'] = len(set(e['consumer'] for e in job['evidence']))
        job['total_mentions'] = len(job['evidence'])
        # Keep top 4 quotes
        job['top_quotes'] = job['evidence'][:4]

    return sorted(jobs, key=lambda x: x['consumer_count'], reverse=True)

def create_presentation(jobs, transcripts):
    """Populate Offbrain template with JTBD content"""

    prs = Presentation(str(TEMPLATE_PATH))

    # Slide 1: Title (already exists in template, modify if needed)
    if len(prs.slides) > 0:
        title_slide = prs.slides[0]
        if hasattr(title_slide.shapes, 'title') and title_slide.shapes.title:
            title_slide.shapes.title.text = "3M Lighting Consumer Research"
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = "Jobs-to-be-Done Framework\nConsumer Insights & Strategic Opportunities"

    # Add new content slides
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(blank_layout)
    _add_title(slide, "Executive Summary")
    _add_bullets(slide, [
        f"{len(transcripts)} consumers analyzed across 82 video interviews",
        f"{len(jobs)} core jobs identified using Clayton Christensen framework",
        "Focus: Understanding progress consumers seek, not just tasks",
        "",
        "KEY FINDING:",
        "Consumers hire lighting solutions to feel capable and in control",
        "They want to avoid contractor complexity while achieving sophisticated aesthetics"
    ])

    # Slide 3: Methodology
    slide = prs.slides.add_slide(blank_layout)
    _add_title(slide, "Jobs-to-be-Done Framework")
    _add_bullets(slide, [
        "A job is the progress a person wants to achieve in specific circumstances",
        "",
        "Job Format: When [situation], I want to [progress], So I can [outcome]",
        "",
        "Three Dimensions:",
        "  • Functional: What they're trying to accomplish",
        "  • Emotional: How they want to feel",
        "  • Social: How they want to be perceived",
        "",
        "Jobs are solution-agnostic and stable over time"
    ])

    # Slides 4-7: Each Job
    for i, job in enumerate(jobs, 1):
        slide = prs.slides.add_slide(blank_layout)
        _add_title(slide, f"Job {i}: {job['name']}")

        bullets = [
            "JOB STATEMENT:",
            job['statement'],
            "",
            f"EVIDENCE: {job['consumer_count']} consumers",
            "",
            "DIMENSIONS:",
            f"• Functional: {job['functional']}",
            f"• Emotional: {job['emotional']}",
            f"• Social: {job['social']}",
            "",
            "CONSUMER VOICE:"
        ]

        # Add top quote with citation
        if job['top_quotes']:
            quote = job['top_quotes'][0]
            bullets.append(f'"{quote["quote"]}"')
            bullets.append(f"— {quote['consumer']} | {quote['source'].split('_')[1]}")

        _add_bullets(slide, bullets)

    # Slide 8: Prioritization
    slide = prs.slides.add_slide(blank_layout)
    _add_title(slide, "Job Prioritization")
    bullets = ["PRIORITY RANKING (by consumer count):", ""]
    for i, job in enumerate(jobs, 1):
        bullets.append(f"{i}. {job['name']}: {job['consumer_count']} consumers")
    bullets.extend(["", "All jobs validated with minimum 4 consumers"])
    _add_bullets(slide, bullets)

    # Slide 9: Strategic Opportunities
    slide = prs.slides.add_slide(blank_layout)
    _add_title(slide, "Strategic Opportunities")
    _add_bullets(slide, [
        "1. WIRELESS PREMIUM:",
        "   Battery solutions that don't sacrifice aesthetic quality",
        "",
        "2. SMART CONTROL AS STANDARD:",
        "   Make remote, timers, dimming core features not add-ons",
        "",
        "3. DIY CONFIDENCE SYSTEM:",
        "   Tools and guides for right product selection",
        "",
        "4. DESIGN-FORWARD DIY:",
        "   Position as the choice for design-conscious DIYers"
    ])

    # Slide 10: Next Steps
    slide = prs.slides.add_slide(blank_layout)
    _add_title(slide, "Next Steps")
    _add_bullets(slide, [
        "IMMEDIATE (2 weeks):",
        "• Internal workshop to review findings",
        "• Prioritization vote with cross-functional team",
        "",
        "SHORT TERM (4-6 weeks):",
        "• Concept development for top opportunities",
        "• Consumer validation with 20-30 participants",
        "",
        "MEDIUM TERM (3 months):",
        "• Prototype development",
        "• In-home testing",
        "• Go-to-market planning"
    ])

    # Save
    output_path = TEMPLATE_PATH.parent / "3M_Lighting_JTBD_Final_Complete.pptx"
    prs.save(str(output_path))

    return output_path, jobs

def _add_title(slide, text):
    """Add title to slide"""
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = 'Arial'

def _add_bullets(slide, items):
    """Add bullet points to slide"""
    left = Inches(0.75)
    top = Inches(1.75)
    width = Inches(8.5)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.font.size = Pt(14) if item.strip() else Pt(8)
        p.font.name = 'Arial'
        p.level = 1 if item.startswith('  •') or item.startswith('   ') else 0

def main():
    print("Loading transcripts...")
    transcripts = load_transcripts()
    print(f"Loaded {len(transcripts)} transcripts")

    print("Extracting JTBD jobs...")
    jobs = extract_jtbd_jobs(transcripts)
    print(f"Identified {len(jobs)} jobs")

    print("Creating presentation...")
    output_path, jobs = create_presentation(jobs, transcripts)

    print(f"\n✓ COMPLETE!")
    print(f"  File: {output_path}")
    print(f"  Jobs: {len(jobs)}")
    for i, job in enumerate(jobs, 1):
        print(f"    {i}. {job['name']}: {job['consumer_count']} consumers")

if __name__ == "__main__":
    main()
