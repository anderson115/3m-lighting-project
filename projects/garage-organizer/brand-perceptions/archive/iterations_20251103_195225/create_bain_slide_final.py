"""
Create a FINAL Bain/McKinsey-quality management consulting slide for 3M Brand Perceptions
Version 3: Perfected consulting-grade design with precise layout control
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

def create_perfect_bain_slide():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # McKinsey/Bain color palette
    NAVY = RGBColor(0, 32, 96)           # Deep navy for headers
    CHARCOAL = RGBColor(64, 64, 64)      # Dark gray for body text
    GRAY = RGBColor(128, 128, 128)       # Medium gray
    LIGHT_GRAY = RGBColor(224, 224, 224) # Light gray for lines
    BLUE = RGBColor(0, 113, 188)         # McKinsey blue
    ORANGE = RGBColor(255, 108, 47)      # Alert orange
    GREEN = RGBColor(0, 176, 80)         # Success green

    # Layout parameters
    MARGIN = Inches(0.5)
    SLIDE_WIDTH = prs.slide_width
    CONTENT_WIDTH = SLIDE_WIDTH - (2 * MARGIN)

    # Top accent line (signature consulting style)
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, Pt(3)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = BLUE
    top_line.line.fill.background()

    # MAIN TITLE
    title_box = slide.shapes.add_textbox(
        MARGIN,
        Inches(0.3),
        CONTENT_WIDTH,
        Inches(0.8)
    )
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Command Brand Analysis: Surface Damage Discussions Outweigh Damage-Free Messaging"
    p.font.name = "Calibri"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # KEY INSIGHT BOX
    insight_top = Inches(1.2)
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN,
        insight_top,
        CONTENT_WIDTH,
        Inches(0.45)
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice blue
    insight_box.line.color.rgb = BLUE
    insight_box.line.width = Pt(1.5)

    insight_text = slide.shapes.add_textbox(
        MARGIN + Inches(0.2),
        insight_top + Inches(0.08),
        CONTENT_WIDTH - Inches(0.4),
        Inches(0.3)
    )
    tf = insight_text.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "KEY FINDING: Surface damage concerns appear in 39% of Command videos (24 of 62), while damage-free benefits are highlighted in only 18% (11 of 62) — a 2.2x gap that challenges core positioning"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = CHARCOAL

    # Create two-column layout
    left_col_start = MARGIN
    left_col_width = Inches(6.0)
    right_col_start = MARGIN + left_col_width + Inches(0.5)
    right_col_width = Inches(6.3)
    content_top = Inches(2.0)

    # LEFT COLUMN - Strategic Analysis

    # Section 1: Context
    context_header = slide.shapes.add_textbox(
        left_col_start,
        content_top,
        left_col_width,
        Inches(0.3)
    )
    tf = context_header.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "SITUATION ASSESSMENT"
    p.font.name = "Calibri"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE

    context_body = slide.shapes.add_textbox(
        left_col_start,
        content_top + Inches(0.3),
        left_col_width,
        Inches(0.8)
    )
    tf = context_body.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Command maintains strong market presence with frequent appearances in creator tutorials. However, content analysis reveals a critical misalignment: creators discuss surface damage issues significantly more than they emphasize damage-free benefits."
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.line_spacing = 1.15

    # Section 2: Implication Box
    impl_top = content_top + Inches(1.3)
    impl_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left_col_start,
        impl_top,
        left_col_width,
        Inches(0.7)
    )
    impl_box.fill.solid()
    impl_box.fill.fore_color.rgb = RGBColor(255, 248, 240)  # Light orange
    impl_box.line.color.rgb = ORANGE
    impl_box.line.width = Pt(2)

    impl_icon = slide.shapes.add_textbox(
        left_col_start + Inches(0.15),
        impl_top + Inches(0.2),
        Inches(0.3),
        Inches(0.3)
    )
    tf = impl_icon.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "!"
    p.font.name = "Calibri"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    impl_text = slide.shapes.add_textbox(
        left_col_start + Inches(0.6),
        impl_top + Inches(0.15),
        left_col_width - Inches(0.8),
        Inches(0.4)
    )
    tf = impl_text.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "BUSINESS IMPACT"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    p = tf.add_paragraph()
    p.text = "When creator feedback contradicts core value proposition, it undermines marketing ROI and erodes brand trust"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL

    # Section 3: Strategic Options
    options_top = impl_top + Inches(0.9)
    options_header = slide.shapes.add_textbox(
        left_col_start,
        options_top,
        left_col_width,
        Inches(0.3)
    )
    tf = options_header.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "STRATEGIC OPTIONS"
    p.font.name = "Calibri"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Options with checkboxes
    options = [
        "Messaging pivot: Emphasize 'removable' over 'damage-free'",
        "Product innovation: Surface-specific adhesive formulations",
        "Education campaign: Clear surface compatibility matrix",
        "Segment focus: Target damage-tolerant applications"
    ]

    option_top = options_top + Inches(0.35)
    for i, option in enumerate(options):
        # Checkbox shape
        checkbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left_col_start + Inches(0.2),
            option_top + (i * Inches(0.35)),
            Inches(0.12),
            Inches(0.12)
        )
        checkbox.fill.solid()
        checkbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        checkbox.line.color.rgb = GRAY
        checkbox.line.width = Pt(1)

        # Option text
        option_text = slide.shapes.add_textbox(
            left_col_start + Inches(0.45),
            option_top + (i * Inches(0.35)) - Inches(0.02),
            left_col_width - Inches(0.5),
            Inches(0.25)
        )
        tf = option_text.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = option
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = CHARCOAL

    # RIGHT COLUMN - Data Visualization

    # Data header
    data_header = slide.shapes.add_textbox(
        right_col_start,
        content_top,
        right_col_width,
        Inches(0.3)
    )
    tf = data_header.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "EVIDENCE FROM CREATOR CONTENT"
    p.font.name = "Calibri"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Chart area
    chart_top = content_top + Inches(0.5)

    # Chart title
    chart_title = slide.shapes.add_textbox(
        right_col_start,
        chart_top,
        right_col_width,
        Inches(0.25)
    )
    tf = chart_title.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Discussion Topic Frequency (n=62 videos)"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.alignment = PP_ALIGN.CENTER

    # Bar chart
    bar_top = chart_top + Inches(0.4)
    bar_data = [
        ("Surface Damage\nDiscussions", 24, ORANGE, "39%"),
        ("Damage-Free\nBenefits", 11, BLUE, "18%")
    ]

    bar_width = Inches(2.0)
    bar_spacing = Inches(1.0)
    max_height = Inches(2.2)
    chart_center = right_col_start + (right_col_width / 2)

    for i, (label, value, color, percent) in enumerate(bar_data):
        bar_left = chart_center - bar_width - (bar_spacing/2) if i == 0 else chart_center + (bar_spacing/2)
        bar_height = (value / 24) * max_height

        # Bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_left,
            bar_top + max_height - bar_height,
            bar_width,
            bar_height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Value on top
        value_text = slide.shapes.add_textbox(
            bar_left,
            bar_top + max_height - bar_height - Inches(0.35),
            bar_width,
            Inches(0.3)
        )
        tf = value_text.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = str(value)
        p.font.name = "Calibri"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Percentage below value
        pct_text = slide.shapes.add_textbox(
            bar_left,
            bar_top + max_height - bar_height - Inches(0.15),
            bar_width,
            Inches(0.2)
        )
        tf = pct_text.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = f"({percent} of videos)"
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        # Label below
        label_text = slide.shapes.add_textbox(
            bar_left - Inches(0.3),
            bar_top + max_height + Inches(0.1),
            bar_width + Inches(0.6),
            Inches(0.35)
        )
        tf = label_text.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = label
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = CHARCOAL
        p.alignment = PP_ALIGN.CENTER

    # Gap indicator
    gap_top = bar_top + max_height + Inches(0.6)
    gap_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_col_start + Inches(1.8),
        gap_top,
        Inches(2.7),
        Inches(0.45)
    )
    gap_box.fill.solid()
    gap_box.fill.fore_color.rgb = RGBColor(255, 245, 230)
    gap_box.line.color.rgb = ORANGE
    gap_box.line.width = Pt(1.5)

    gap_text = slide.shapes.add_textbox(
        right_col_start + Inches(1.9),
        gap_top + Inches(0.08),
        Inches(2.5),
        Inches(0.3)
    )
    tf = gap_text.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "2.2x"
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Discussion Gap"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.alignment = PP_ALIGN.CENTER

    # Sentiment bar
    sent_top = gap_top + Inches(0.7)
    sent_title = slide.shapes.add_textbox(
        right_col_start,
        sent_top,
        right_col_width,
        Inches(0.25)
    )
    tf = sent_title.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Overall Sentiment Distribution"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.alignment = PP_ALIGN.CENTER

    # Sentiment segments
    sent_bar_top = sent_top + Inches(0.3)
    sent_bar_left = right_col_start + Inches(0.5)
    sent_bar_width = Inches(5.3)
    sent_bar_height = Inches(0.35)

    segments = [
        (0.84, GRAY, "Neutral", "84%"),
        (0.07, GREEN, "Positive", "7%"),
        (0.09, ORANGE, "Negative", "9%")
    ]

    current_left = sent_bar_left
    for proportion, color, label, value in segments:
        seg_width = sent_bar_width * proportion

        # Segment bar
        segment = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            current_left,
            sent_bar_top,
            seg_width,
            sent_bar_height
        )
        segment.fill.solid()
        segment.fill.fore_color.rgb = color
        segment.line.fill.background()

        # Label for large segments
        if proportion > 0.2:
            seg_text = slide.shapes.add_textbox(
                current_left,
                sent_bar_top + Inches(0.08),
                seg_width,
                Inches(0.2)
            )
            tf = seg_text.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = f"{label} {value}"
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        current_left += seg_width

    # Small segment labels
    small_labels = slide.shapes.add_textbox(
        sent_bar_left + (sent_bar_width * 0.84),
        sent_bar_top + sent_bar_height + Inches(0.05),
        Inches(2.0),
        Inches(0.2)
    )
    tf = small_labels.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Positive 7% | Negative 9%"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY

    # BOTTOM SECTION - Next Steps
    bottom_top = Inches(5.8)

    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN,
        bottom_top,
        CONTENT_WIDTH,
        Pt(1)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = LIGHT_GRAY
    divider.line.fill.background()

    # Next steps box
    next_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN,
        bottom_top + Inches(0.15),
        CONTENT_WIDTH,
        Inches(0.65)
    )
    next_box.fill.solid()
    next_box.fill.fore_color.rgb = RGBColor(248, 248, 252)
    next_box.line.fill.background()

    # Next steps header
    next_header = slide.shapes.add_textbox(
        MARGIN + Inches(0.3),
        bottom_top + Inches(0.2),
        Inches(2.0),
        Inches(0.25)
    )
    tf = next_header.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "KEY QUESTIONS"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE

    # Questions
    questions = [
        "How might we better align damage-free positioning with creator experiences?",
        "Should investment prioritize product reformulation or education initiatives?",
        "What role should Command play in segments where damage-free is unachievable?"
    ]

    q_left = MARGIN + Inches(2.8)
    q_width = Inches(3.5)

    for i, q in enumerate(questions):
        q_text = slide.shapes.add_textbox(
            q_left + (i * Inches(3.6)),
            bottom_top + Inches(0.25),
            q_width,
            Inches(0.4)
        )
        tf = q_text.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = f"{i+1}. {q}"
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.color.rgb = CHARCOAL
        p.line_spacing = 1.1

    # FOOTER
    footer_top = Inches(6.9)

    # Footer line
    footer_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN,
        footer_top,
        CONTENT_WIDTH,
        Pt(0.5)
    )
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = LIGHT_GRAY
    footer_line.line.fill.background()

    # Source
    source = slide.shapes.add_textbox(
        MARGIN,
        footer_top + Inches(0.1),
        Inches(5.0),
        Inches(0.2)
    )
    tf = source.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Source: Phase 2 Analysis of YouTube creator content (62 Command-focused videos)"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY

    # Date and page
    date_page = slide.shapes.add_textbox(
        SLIDE_WIDTH - MARGIN - Inches(2.5),
        footer_top + Inches(0.1),
        Inches(2.5),
        Inches(0.2)
    )
    tf = date_page.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "3M Brand Perceptions | November 2024 | 1"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

    # Save
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"/Users/anderson115/00-interlink/12-work/3m-lighting-project/modules/brand-perceptions/3M_BrandPerceptions_BainQuality_FINAL_{timestamp}.pptx"
    prs.save(filename)

    return filename

if __name__ == "__main__":
    filename = create_perfect_bain_slide()
    print(f"Final Bain-quality slide created: {filename}")