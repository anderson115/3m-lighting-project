"""
Create a TRUE Bain/McKinsey-quality management consulting slide for 3M Brand Perceptions
Version 2: Professional consulting-grade design with proper layout control
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import datetime

def add_text_box(slide, left, top, width, height, text, font_size, bold=False,
                 color=None, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Helper function to add formatted text boxes"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.margin_left = 0
    text_frame.margin_right = 0

    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment

    return textbox

def create_professional_bain_slide():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Standard widescreen
    prs.slide_height = Inches(7.5)

    # Add completely blank slide (no layout)
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Set slide background to pure white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Professional color palette (Bain/McKinsey style)
    CHARCOAL = RGBColor(51, 51, 51)      # Main text
    DARK_GRAY = RGBColor(102, 102, 102)   # Secondary text
    MED_GRAY = RGBColor(153, 153, 153)    # Supporting text
    LIGHT_GRAY = RGBColor(217, 217, 217)  # Borders/lines
    BAIN_BLUE = RGBColor(0, 84, 166)      # Primary accent
    ACCENT_ORANGE = RGBColor(255, 125, 0)  # Secondary accent
    ACCENT_GREEN = RGBColor(0, 166, 81)    # Positive

    # Professional margins
    LEFT_MARGIN = Inches(0.75)
    RIGHT_MARGIN = Inches(0.75)
    TOP_MARGIN = Inches(0.5)
    CONTENT_WIDTH = prs.slide_width - LEFT_MARGIN - RIGHT_MARGIN

    # PROFESSIONAL HEADER BAR
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        Inches(0.08)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = BAIN_BLUE
    header_bar.line.fill.background()

    # MAIN HEADLINE - Professional styling
    headline_text = ("Command's Market Presence is Strong, But Creator "
                    "Discussions Reveal Strategic Tension")

    add_text_box(
        slide,
        LEFT_MARGIN,
        Inches(0.4),
        CONTENT_WIDTH,
        Inches(0.6),
        headline_text,
        Pt(28),
        bold=True,
        color=CHARCOAL,
        font_name="Calibri"
    )

    # EXECUTIVE SUMMARY BOX - Key takeaway upfront
    summary_top = Inches(1.2)
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LEFT_MARGIN,
        summary_top,
        CONTENT_WIDTH,
        Inches(0.5)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(245, 247, 250)  # Very light blue-gray
    summary_box.line.color.rgb = BAIN_BLUE
    summary_box.line.width = Pt(1)

    summary_text = ("KEY INSIGHT: Surface damage concerns appear 2.2x more frequently than "
                   "damage-free benefits in creator content, suggesting expectation misalignment")

    add_text_box(
        slide,
        LEFT_MARGIN + Inches(0.2),
        summary_top + Inches(0.1),
        CONTENT_WIDTH - Inches(0.4),
        Inches(0.3),
        summary_text,
        Pt(13),
        bold=False,
        color=CHARCOAL,
        font_name="Calibri"
    )

    # MAIN CONTENT AREA - Two columns
    content_top = Inches(2.0)
    left_column_width = Inches(5.5)
    right_column_width = Inches(6.0)
    column_gap = Inches(0.5)

    # LEFT COLUMN - Analysis & Insights

    # Section 1: The Pattern
    pattern_header = add_text_box(
        slide,
        LEFT_MARGIN,
        content_top,
        left_column_width,
        Inches(0.3),
        "OBSERVED PATTERN",
        Pt(14),
        bold=True,
        color=BAIN_BLUE,
        font_name="Calibri"
    )

    pattern_text = ("Command appears frequently in creator tutorials, particularly for "
                   "picture hanging. However, analysis reveals that surface damage "
                   "discussions significantly outweigh damage-free benefit mentions, "
                   "creating potential brand promise challenges.")

    add_text_box(
        slide,
        LEFT_MARGIN,
        content_top + Inches(0.35),
        left_column_width,
        Inches(0.7),
        pattern_text,
        Pt(12),
        color=CHARCOAL,
        font_name="Calibri"
    )

    # Section 2: Key Finding Box
    finding_top = content_top + Inches(1.2)
    finding_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LEFT_MARGIN,
        finding_top,
        left_column_width,
        Inches(0.8)
    )
    finding_box.fill.solid()
    finding_box.fill.fore_color.rgb = RGBColor(255, 250, 245)  # Very light orange
    finding_box.line.color.rgb = ACCENT_ORANGE
    finding_box.line.width = Pt(2)

    add_text_box(
        slide,
        LEFT_MARGIN + Inches(0.2),
        finding_top + Inches(0.1),
        left_column_width - Inches(0.4),
        Inches(0.15),
        "CRITICAL FINDING",
        Pt(11),
        bold=True,
        color=ACCENT_ORANGE,
        font_name="Calibri"
    )

    finding_detail = ("When creators discuss surface damage (39% of videos), "
                     "it directly contradicts Command's core value proposition")

    add_text_box(
        slide,
        LEFT_MARGIN + Inches(0.2),
        finding_top + Inches(0.3),
        left_column_width - Inches(0.4),
        Inches(0.4),
        finding_detail,
        Pt(13),
        bold=True,
        color=CHARCOAL,
        font_name="Calibri"
    )

    # Section 3: Strategic Options
    options_top = finding_top + Inches(1.0)

    add_text_box(
        slide,
        LEFT_MARGIN,
        options_top,
        left_column_width,
        Inches(0.3),
        "STRATEGIC CONSIDERATIONS",
        Pt(14),
        bold=True,
        color=BAIN_BLUE,
        font_name="Calibri"
    )

    # Create professional bullet points
    bullet_points = [
        "• Reframe messaging: Focus on 'removable' vs. 'damage-free'",
        "• Product innovation: Develop surface-specific formulations",
        "• Education initiative: Clear compatibility guidelines",
        "• Market segmentation: Target damage-tolerant segments"
    ]

    bullet_top = options_top + Inches(0.35)
    for i, bullet in enumerate(bullet_points):
        add_text_box(
            slide,
            LEFT_MARGIN + Inches(0.2),
            bullet_top + (i * Inches(0.35)),
            left_column_width - Inches(0.4),
            Inches(0.3),
            bullet,
            Pt(12),
            color=CHARCOAL,
            font_name="Calibri"
        )

    # RIGHT COLUMN - Data Visualization

    # Title for data section
    add_text_box(
        slide,
        LEFT_MARGIN + left_column_width + column_gap,
        content_top,
        right_column_width,
        Inches(0.3),
        "SUPPORTING EVIDENCE",
        Pt(14),
        bold=True,
        color=BAIN_BLUE,
        font_name="Calibri"
    )

    # Data Visualization 1: Comparison Chart
    chart_left = LEFT_MARGIN + left_column_width + column_gap
    chart_top = content_top + Inches(0.4)

    # Create visual comparison using shapes (more control than chart)

    # Chart title
    add_text_box(
        slide,
        chart_left,
        chart_top,
        right_column_width,
        Inches(0.25),
        "Creator Discussion Frequency",
        Pt(12),
        bold=True,
        color=CHARCOAL,
        font_name="Calibri",
        alignment=PP_ALIGN.CENTER
    )

    # Create bar chart visualization
    bar_data = [
        ("Surface Damage\nDiscussions", 24, ACCENT_ORANGE),
        ("Damage-Free\nBenefit Mentions", 11, BAIN_BLUE)
    ]

    bar_top = chart_top + Inches(0.4)
    bar_width = Inches(1.8)
    bar_spacing = Inches(2.5)
    max_height = Inches(2.0)

    for i, (label, value, color) in enumerate(bar_data):
        bar_left = chart_left + Inches(0.5) + (i * bar_spacing)
        bar_height = (value / 24) * max_height  # Normalize to max value

        # Draw bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_left,
            bar_top + max_height - bar_height,
            bar_width,
            bar_height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Add value label on top of bar
        add_text_box(
            slide,
            bar_left,
            bar_top + max_height - bar_height - Inches(0.3),
            bar_width,
            Inches(0.25),
            str(value),
            Pt(20),
            bold=True,
            color=color,
            font_name="Calibri",
            alignment=PP_ALIGN.CENTER
        )

        # Add label below bar
        add_text_box(
            slide,
            bar_left - Inches(0.2),
            bar_top + max_height + Inches(0.1),
            bar_width + Inches(0.4),
            Inches(0.4),
            label,
            Pt(11),
            bold=False,
            color=CHARCOAL,
            font_name="Calibri",
            alignment=PP_ALIGN.CENTER
        )

    # Add ratio indicator
    ratio_top = bar_top + max_height + Inches(0.7)
    ratio_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        chart_left + Inches(1.5),
        ratio_top,
        Inches(3.0),
        Inches(0.5)
    )
    ratio_box.fill.solid()
    ratio_box.fill.fore_color.rgb = RGBColor(255, 245, 235)
    ratio_box.line.color.rgb = ACCENT_ORANGE
    ratio_box.line.width = Pt(1)

    add_text_box(
        slide,
        chart_left + Inches(1.7),
        ratio_top + Inches(0.12),
        Inches(2.6),
        Inches(0.25),
        "2.2x Gap in Discussion",
        Pt(14),
        bold=True,
        color=ACCENT_ORANGE,
        font_name="Calibri",
        alignment=PP_ALIGN.CENTER
    )

    # Sentiment Distribution
    sentiment_top = ratio_top + Inches(0.8)

    add_text_box(
        slide,
        chart_left,
        sentiment_top,
        right_column_width,
        Inches(0.25),
        "Overall Sentiment Distribution",
        Pt(12),
        bold=True,
        color=CHARCOAL,
        font_name="Calibri",
        alignment=PP_ALIGN.CENTER
    )

    # Create sentiment bar
    sent_bar_top = sentiment_top + Inches(0.3)
    sent_bar_left = chart_left + Inches(0.5)
    sent_bar_width = Inches(5.0)
    sent_bar_height = Inches(0.3)

    # Segments
    segments = [
        (0.84, MED_GRAY, "Neutral 84%"),
        (0.065, ACCENT_GREEN, "Pos 7%"),
        (0.065, ACCENT_ORANGE, "Neg 7%")
    ]

    current_left = sent_bar_left
    for proportion, color, label in segments:
        segment_width = sent_bar_width * proportion

        segment = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            current_left,
            sent_bar_top,
            segment_width,
            sent_bar_height
        )
        segment.fill.solid()
        segment.fill.fore_color.rgb = color
        segment.line.fill.background()

        # Add label if segment is large enough
        if proportion > 0.15:
            add_text_box(
                slide,
                current_left,
                sent_bar_top + Inches(0.05),
                segment_width,
                Inches(0.2),
                label,
                Pt(10),
                bold=True,
                color=RGBColor(255, 255, 255),
                font_name="Calibri",
                alignment=PP_ALIGN.CENTER
            )

        current_left += segment_width

    # Add small segment labels below
    add_text_box(
        slide,
        sent_bar_left + sent_bar_width * 0.84,
        sent_bar_top + sent_bar_height + Inches(0.05),
        Inches(1.0),
        Inches(0.2),
        "Pos 7% | Neg 7%",
        Pt(9),
        color=DARK_GRAY,
        font_name="Calibri"
    )

    # BOTTOM ACTION BOX - Call to action
    action_top = Inches(6.3)
    action_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LEFT_MARGIN,
        action_top,
        CONTENT_WIDTH,
        Inches(0.6)
    )
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(240, 240, 245)
    action_box.line.fill.background()

    add_text_box(
        slide,
        LEFT_MARGIN + Inches(0.3),
        action_top + Inches(0.08),
        Inches(2.0),
        Inches(0.4),
        "KEY QUESTIONS:",
        Pt(12),
        bold=True,
        color=BAIN_BLUE,
        font_name="Calibri"
    )

    questions = [
        "1. How might we align damage-free claims with actual creator experiences?",
        "2. Should we invest in surface education or product reformulation?",
        "3. What role should Command play if damage-free isn't always achievable?"
    ]

    question_left = LEFT_MARGIN + Inches(2.5)
    for i, q in enumerate(questions):
        add_text_box(
            slide,
            question_left + (i * Inches(3.5)),
            action_top + Inches(0.15),
            Inches(3.2),
            Inches(0.4),
            q,
            Pt(11),
            color=CHARCOAL,
            font_name="Calibri"
        )

    # FOOTER - Professional finishing
    footer_top = Inches(7.1)

    # Subtle divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        LEFT_MARGIN,
        footer_top,
        CONTENT_WIDTH,
        Pt(0.75)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = LIGHT_GRAY
    divider.line.fill.background()

    # Source and page number
    add_text_box(
        slide,
        LEFT_MARGIN,
        footer_top + Inches(0.1),
        Inches(4),
        Inches(0.2),
        "Source: Phase 2 Analysis of 62 Command creator videos",
        Pt(9),
        color=MED_GRAY,
        font_name="Calibri"
    )

    add_text_box(
        slide,
        prs.slide_width - RIGHT_MARGIN - Inches(1.5),
        footer_top + Inches(0.1),
        Inches(1.5),
        Inches(0.2),
        "3M Brand Perceptions | 1",
        Pt(9),
        color=MED_GRAY,
        font_name="Calibri",
        alignment=PP_ALIGN.RIGHT
    )

    # Save presentation
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"/Users/anderson115/00-interlink/12-work/3m-lighting-project/modules/brand-perceptions/3M_BrandPerceptions_BainQuality_v2_{timestamp}.pptx"
    prs.save(filename)

    return filename

if __name__ == "__main__":
    filename = create_professional_bain_slide()
    print(f"Professional Bain-quality slide created: {filename}")