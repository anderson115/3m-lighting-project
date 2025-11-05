from pptx import Presentation
from pptx.util import Inches

prs = Presentation('SLIDE2_pythonpptx.pptx')
slide = prs.slides[0]

# Adjust third performance bar to increase width and align value text
for shape in slide.shapes:
    if shape.name == 'Rounded Rectangle 23':
        shape.width = Inches(6.2)
    if shape.name == 'TextBox 24':
        shape.left = shape.left - Inches(2.5)

prs.save('SLIDE2_pythonpptx.pptx')
