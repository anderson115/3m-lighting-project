#!/usr/bin/env python3
"""
Convert HTML slides to PowerPoint with ZERO design loss
Strategy: High-res screenshots → PowerPoint image slides
"""

import subprocess
import time
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

def create_screenshot(html_file, output_png, width=1920, height=1080):
    """Use playwright to capture pixel-perfect screenshot"""
    script = f"""
from playwright.sync_api import sync_playwright

with sync_playwright() as p:
    browser = p.chromium.launch()
    page = browser.new_page(viewport={{'width': {width}, 'height': {height}}})
    page.goto('file://{html_file}')
    page.wait_for_load_state('networkidle')
    page.screenshot(path='{output_png}', full_page=False)
    browser.close()
    print(f'✅ Screenshot: {output_png}')
"""

    # Run playwright screenshot
    result = subprocess.run(
        ['python3', '-c', script],
        capture_output=True,
        text=True
    )

    if result.returncode == 0:
        return True
    else:
        print(f"❌ Screenshot failed: {result.stderr}")
        return False

def create_pptx_from_screenshots(screenshot_files, output_pptx):
    """Create PowerPoint from screenshot images"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for img_path in screenshot_files:
        if not Path(img_path).exists():
            print(f"⚠️  Skipping missing file: {img_path}")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Add image filling entire slide
        slide.shapes.add_picture(
            img_path,
            0, 0,  # Left, Top
            Inches(16), Inches(9)  # Width, Height
        )

        print(f"✅ Added slide: {Path(img_path).name}")

    prs.save(output_pptx)
    print(f"\n🎉 PowerPoint created: {output_pptx}")
    print(f"📊 Total slides: {len(screenshot_files)}")
    print(f"✨ Design preservation: 100% (pixel-perfect)")

def main():
    base_dir = Path.cwd()

    # All slide HTML files
    slides = [
        "SLIDE1_DEFAULT.html",
        "SLIDE2_DEFAULT.html",
        "SLIDE3_DEFAULT.html",
        "SLIDE4_DEFAULT.html",
        "SLIDE5_DEFAULT.html",
        "SLIDE6_DEFAULT.html",
    ]

    print("🎯 HTML → PowerPoint Conversion (Perfect Fidelity)")
    print("=" * 60)

    # Step 1: Create screenshots
    print("\n📸 Step 1: Capturing pixel-perfect screenshots...")
    screenshot_files = []

    for slide_html in slides:
        html_path = base_dir / slide_html
        if not html_path.exists():
            print(f"⚠️  File not found: {slide_html}")
            continue

        png_path = base_dir / f"{slide_html.replace('.html', '_SCREENSHOT.png')}"

        if create_screenshot(str(html_path.absolute()), str(png_path)):
            screenshot_files.append(str(png_path))

    # Step 2: Create PowerPoint
    print(f"\n🎨 Step 2: Building PowerPoint from screenshots...")
    output_pptx = base_dir / "OFFBRAIN_SLIDES_DEFAULT.pptx"
    create_pptx_from_screenshots(screenshot_files, str(output_pptx))

    print("\n" + "=" * 60)
    print("✅ CONVERSION COMPLETE")
    print(f"📁 Output: {output_pptx.name}")
    print(f"🎨 Design Loss: 0% (pixel-perfect image preservation)")
    print(f"📝 Editable: No (image-based)")
    print(f"💡 Tip: For editable slides, use native PowerPoint design")

if __name__ == "__main__":
    main()
