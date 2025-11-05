"""
Create a Bain-quality management consulting slide for 3M Brand Perceptions
Based on detailed design specifications for Slide 1: Command Brand Perception
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import datetime

def create_bain_quality_slide():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Add blank slide
    blank_slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Define colors (matching Bain/McKinsey palette)
    DEEP_GRAY = RGBColor(17, 24, 39)
    SLATE_GRAY = RGBColor(55, 65, 81)
    MEDIUM_GRAY = RGBColor(107, 114, 128)
    LIGHT_GRAY = RGBColor(156, 163, 175)
    BLUE_PRIMARY = RGBColor(30, 58, 138)
    AMBER_ACCENT = RGBColor(245, 158, 11)
    TEAL_SUPPORT = RGBColor(8, 145, 178)
    PANEL_BG = RGBColor(249, 250, 251)

    # Margins (converted from pixels to inches - 80px = 0.833 inches at 96 DPI)
    margin = Inches(0.833)

    # HEADLINE
    headline = slide.shapes.add_textbox(
        margin,
        margin,
        Inches(8.334),  # Full width minus margins
        Inches(0.8)
    )
    headline_frame = headline.text_frame
    headline_frame.clear()
    headline_frame.word_wrap = True
    headline_frame.margin_top = 0
    headline_frame.margin_bottom = 0
    headline_frame.margin_left = 0
    headline_frame.margin_right = 0

    p = headline_frame.add_paragraph()
    p.text = "Command's Market Presence is Strong, But Creator\nDiscussions Reveal a Tension Worth Exploring"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DEEP_GRAY
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.1875  # 38pt/32pt

    # LEFT COLUMN CONTENT
    left_col_top = Inches(1.8)
    left_col_width = Inches(3.5)

    # THE PATTERN Section
    pattern_box = slide.shapes.add_textbox(
        margin,
        left_col_top,
        left_col_width,
        Inches(1.2)
    )
    pattern_frame = pattern_box.text_frame
    pattern_frame.clear()
    pattern_frame.word_wrap = True
    pattern_frame.margin_top = 0
    pattern_frame.margin_bottom = 0
    pattern_frame.margin_left = 0
    pattern_frame.margin_right = 0

    # Add asterisk and header
    p = pattern_frame.add_paragraph()
    p.text = "✱ THE PATTERN"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"
    p.space_after = Pt(8)

    # Add narrative text
    p = pattern_frame.add_paragraph()
    p.text = ("Command appears frequently in creator content—particularly in tutorials. "
              "However, surface damage appears in feedback more often than damage-free "
              "benefits are highlighted, creating a strategic question about expectation setting.")
    p.font.size = Pt(14)
    p.font.color.rgb = DEEP_GRAY
    p.font.name = "Helvetica Neue"
    p.line_spacing = 1.43  # 20pt/14pt

    # HARD TRUTH Panel
    panel_top = Inches(2.8)
    panel_height = Inches(0.9)

    # Create panel background
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        margin,
        panel_top,
        left_col_width,
        panel_height
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL_BG
    panel.line.color.rgb = BLUE_PRIMARY
    panel.line.width = Pt(3)

    # Only show left border - hide others
    panel.line.fill.background()

    # Add blue diamond (as text for simplicity)
    diamond_box = slide.shapes.add_textbox(
        margin + Inches(0.15),
        panel_top + Inches(0.15),
        Inches(0.3),
        Inches(0.3)
    )
    diamond_frame = diamond_box.text_frame
    diamond_frame.clear()
    p = diamond_frame.add_paragraph()
    p.text = "🔷"
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE_PRIMARY

    # Add Hard Truth text
    truth_box = slide.shapes.add_textbox(
        margin + Inches(0.5),
        panel_top + Inches(0.15),
        left_col_width - Inches(0.7),
        panel_height - Inches(0.3)
    )
    truth_frame = truth_box.text_frame
    truth_frame.clear()
    truth_frame.word_wrap = True
    truth_frame.margin_top = 0
    truth_frame.margin_bottom = 0

    p = truth_frame.add_paragraph()
    p.text = "HARD TRUTH"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"

    p = truth_frame.add_paragraph()
    p.text = "Surface damage concerns outweigh damage-free benefit emphasis in creator content"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DEEP_GRAY
    p.font.name = "Helvetica Neue"
    p.space_before = Pt(4)

    p = truth_frame.add_paragraph()
    p.text = "This matters because marketing ROI is compromised when creator feedback contradicts positioning"
    p.font.size = Pt(12)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = "Helvetica Neue"
    p.space_before = Pt(6)

    # PROVOCATION Section
    provocation_top = Inches(3.9)

    provocation_box = slide.shapes.add_textbox(
        margin,
        provocation_top,
        left_col_width,
        Inches(0.8)
    )
    provocation_frame = provocation_box.text_frame
    provocation_frame.clear()
    provocation_frame.word_wrap = True
    provocation_frame.margin_top = 0

    # Add question with icon
    p = provocation_frame.add_paragraph()
    p.text = "❓ PROVOCATION"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT
    p.font.name = "Helvetica Neue"
    p.space_after = Pt(6)

    p = provocation_frame.add_paragraph()
    p.text = '"Are claims setting expectations too high?"'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DEEP_GRAY
    p.font.name = "Helvetica Neue"
    p.space_after = Pt(8)

    p = provocation_frame.add_paragraph()
    p.text = "Consider:"
    p.font.size = Pt(13)
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"
    p.space_after = Pt(4)

    # Add bullet points
    options = [
        "• Product innovation to reduce damage",
        "• Transparent compatibility labeling",
        "• Reframe expectations in messaging"
    ]

    for option in options:
        p = provocation_frame.add_paragraph()
        p.text = option
        p.font.size = Pt(13)
        p.font.color.rgb = DEEP_GRAY
        p.font.name = "Helvetica Neue"
        p.level = 0
        p.space_after = Pt(2)

    # RIGHT COLUMN CONTENT
    right_col_left = Inches(4.7)
    right_col_width = Inches(4.6)

    # Balance Scale Visual (using shapes)
    scale_center_x = right_col_left + right_col_width/2
    scale_top = left_col_top + Inches(0.2)

    # Fulcrum (triangle)
    fulcrum = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        scale_center_x - Inches(0.15),
        scale_top + Inches(1.1),
        Inches(0.3),
        Inches(0.3)
    )
    fulcrum.fill.solid()
    fulcrum.fill.fore_color.rgb = SLATE_GRAY
    fulcrum.line.fill.background()

    # Beam (rectangle, tilted)
    beam_width = Inches(3.0)
    beam = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        scale_center_x - beam_width/2,
        scale_top + Inches(0.9),
        beam_width,
        Pt(3)
    )
    beam.fill.solid()
    beam.fill.fore_color.rgb = MEDIUM_GRAY
    beam.line.fill.background()
    beam.rotation = -8  # Tilt to show imbalance

    # Left platform (heavier/lower)
    left_platform = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        scale_center_x - beam_width/2 - Inches(0.2),
        scale_top + Inches(0.7),
        Inches(0.8),
        Inches(0.5)
    )
    left_platform.fill.solid()
    left_platform.fill.fore_color.rgb = AMBER_ACCENT
    left_platform.line.color.rgb = AMBER_ACCENT
    left_platform.line.width = Pt(2)

    # Right platform (lighter/higher)
    right_platform = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        scale_center_x + beam_width/2 - Inches(0.6),
        scale_top + Inches(0.5),
        Inches(0.8),
        Inches(0.5)
    )
    right_platform.fill.solid()
    right_platform.fill.fore_color.rgb = BLUE_PRIMARY
    right_platform.line.color.rgb = BLUE_PRIMARY
    right_platform.line.width = Pt(2)

    # Labels for platforms
    left_label = slide.shapes.add_textbox(
        scale_center_x - beam_width/2 - Inches(0.3),
        scale_top + Inches(1.3),
        Inches(1.0),
        Inches(0.3)
    )
    left_label_frame = left_label.text_frame
    left_label_frame.clear()
    p = left_label_frame.add_paragraph()
    p.text = "Surface Damage\nDiscussions"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.CENTER

    right_label = slide.shapes.add_textbox(
        scale_center_x + beam_width/2 - Inches(0.6),
        scale_top + Inches(1.3),
        Inches(1.0),
        Inches(0.3)
    )
    right_label_frame = right_label.text_frame
    right_label_frame.clear()
    p = right_label_frame.add_paragraph()
    p.text = "Damage-Free\nBenefits"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.CENTER

    # Data Table below scale
    table_top = scale_top + Inches(1.8)
    table_left = right_col_left + Inches(0.5)
    table_width = Inches(3.5)

    # Table header
    table_header = slide.shapes.add_textbox(
        table_left,
        table_top,
        table_width,
        Inches(0.25)
    )
    header_frame = table_header.text_frame
    header_frame.clear()
    p = header_frame.add_paragraph()
    p.text = "DISCUSSION ANALYSIS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.LEFT

    # Create table using shapes and text boxes
    rows_data = [
        ("Discussion Type", "Videos"),
        ("Surface Damage", "24"),
        ("Damage-Free Emphasis", "11")
    ]

    row_height = Inches(0.3)
    current_top = table_top + Inches(0.3)

    for i, (label, value) in enumerate(rows_data):
        # Row background
        if i == 0:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                table_left,
                current_top,
                table_width,
                row_height
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = PANEL_BG
            row_bg.line.color.rgb = BLUE_PRIMARY
            row_bg.line.width = Pt(2) if i == 0 else Pt(1)

        # Label text
        label_box = slide.shapes.add_textbox(
            table_left + Inches(0.1),
            current_top,
            table_width * 0.7,
            row_height
        )
        label_frame = label_box.text_frame
        label_frame.clear()
        label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = label_frame.add_paragraph()
        p.text = label
        p.font.size = Pt(12) if i == 0 else Pt(14)
        p.font.bold = i == 0
        p.font.color.rgb = SLATE_GRAY if i == 0 else DEEP_GRAY
        p.font.name = "Helvetica Neue"

        # Value text
        value_box = slide.shapes.add_textbox(
            table_left + table_width * 0.7,
            current_top,
            table_width * 0.25,
            row_height
        )
        value_frame = value_box.text_frame
        value_frame.clear()
        value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = value_frame.add_paragraph()
        p.text = value
        p.font.size = Pt(12) if i == 0 else Pt(14)
        p.font.bold = True
        p.font.color.rgb = SLATE_GRAY if i == 0 else DEEP_GRAY
        p.font.name = "Helvetica Neue"
        p.alignment = PP_ALIGN.RIGHT

        current_top += row_height

    # Sentiment Distribution Bar
    sentiment_top = current_top + Inches(0.3)

    # Title
    sentiment_title = slide.shapes.add_textbox(
        table_left,
        sentiment_top,
        table_width,
        Inches(0.25)
    )
    title_frame = sentiment_title.text_frame
    title_frame.clear()
    p = title_frame.add_paragraph()
    p.text = "SENTIMENT DISTRIBUTION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"

    # Create stacked bar
    bar_top = sentiment_top + Inches(0.3)
    bar_height = Inches(0.35)

    # Neutral segment (84%)
    neutral_width = table_width * 0.84
    neutral_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        table_left,
        bar_top,
        neutral_width,
        bar_height
    )
    neutral_bar.fill.solid()
    neutral_bar.fill.fore_color.rgb = MEDIUM_GRAY
    neutral_bar.line.fill.background()

    # Add text on neutral bar
    neutral_text = slide.shapes.add_textbox(
        table_left + neutral_width/2 - Inches(0.4),
        bar_top + Inches(0.05),
        Inches(0.8),
        Inches(0.25)
    )
    neutral_frame = neutral_text.text_frame
    neutral_frame.clear()
    p = neutral_frame.add_paragraph()
    p.text = "Neutral 84%"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.CENTER

    # Positive segment (6.5%)
    positive_width = table_width * 0.065
    positive_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        table_left + neutral_width,
        bar_top,
        positive_width,
        bar_height
    )
    positive_bar.fill.solid()
    positive_bar.fill.fore_color.rgb = BLUE_PRIMARY
    positive_bar.line.fill.background()

    # Negative segment (6.5%)
    negative_width = table_width * 0.065
    negative_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        table_left + neutral_width + positive_width,
        bar_top,
        negative_width,
        bar_height
    )
    negative_bar.fill.solid()
    negative_bar.fill.fore_color.rgb = AMBER_ACCENT
    negative_bar.line.fill.background()

    # Add labels for small segments below bar
    labels_top = bar_top + bar_height + Inches(0.05)

    pos_label = slide.shapes.add_textbox(
        table_left + neutral_width - Inches(0.2),
        labels_top,
        Inches(0.7),
        Inches(0.2)
    )
    pos_frame = pos_label.text_frame
    pos_frame.clear()
    p = pos_frame.add_paragraph()
    p.text = "Pos 6.5%"
    p.font.size = Pt(10)
    p.font.color.rgb = BLUE_PRIMARY
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.CENTER

    neg_label = slide.shapes.add_textbox(
        table_left + neutral_width + positive_width,
        labels_top,
        Inches(0.7),
        Inches(0.2)
    )
    neg_frame = neg_label.text_frame
    neg_frame.clear()
    p = neg_frame.add_paragraph()
    p.text = "Neg 6.5%"
    p.font.size = Pt(10)
    p.font.color.rgb = AMBER_ACCENT
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.CENTER

    # CONVERSATION STARTERS at bottom
    conv_top = Inches(4.8)

    conv_box = slide.shapes.add_textbox(
        margin,
        conv_top,
        Inches(8.334),
        Inches(0.4)
    )
    conv_frame = conv_box.text_frame
    conv_frame.clear()

    p = conv_frame.add_paragraph()
    p.text = "→ CONVERSATION STARTERS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SLATE_GRAY
    p.font.name = "Helvetica Neue"
    p.space_after = Pt(4)

    questions = [
        "→ How might we align damage-free claims with actual creator experiences?",
        "→ Should we invest in surface compatibility education or product reformulation?",
        "→ What role should Command play if damage-free isn't always achievable?"
    ]

    for q in questions:
        p = conv_frame.add_paragraph()
        p.text = q
        p.font.size = Pt(12)
        p.font.color.rgb = DEEP_GRAY
        p.font.name = "Helvetica Neue"
        p.level = 0
        p.space_after = Pt(2)

    # FOOTER
    footer_top = Inches(5.2)

    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        margin,
        footer_top,
        Inches(8.334),
        Pt(1)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(229, 231, 235)
    divider.line.fill.background()

    # Source citation
    source_box = slide.shapes.add_textbox(
        margin,
        footer_top + Inches(0.1),
        Inches(4),
        Inches(0.2)
    )
    source_frame = source_box.text_frame
    source_frame.clear()
    p = source_frame.add_paragraph()
    p.text = "Source: Phase 2 Analysis, 62 Command videos"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Helvetica Neue"

    # Page number
    page_box = slide.shapes.add_textbox(
        Inches(8.5),
        footer_top + Inches(0.1),
        Inches(1),
        Inches(0.2)
    )
    page_frame = page_box.text_frame
    page_frame.clear()
    p = page_frame.add_paragraph()
    p.text = "Slide 1 of 3"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Helvetica Neue"
    p.alignment = PP_ALIGN.RIGHT

    # Save the presentation
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"/Users/anderson115/00-interlink/12-work/3m-lighting-project/modules/brand-perceptions/3M_BrandPerceptions_BainQuality_{timestamp}.pptx"
    prs.save(filename)

    return filename

if __name__ == "__main__":
    filename = create_bain_quality_slide()
    print(f"Bain-quality slide created: {filename}")