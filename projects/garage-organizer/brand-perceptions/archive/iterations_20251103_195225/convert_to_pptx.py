#!/usr/bin/env python3
"""
Convert offbrain HTML slides to PowerPoint (.pptx) format
Preserves exact design specifications from HTML/CSS
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Design System Colors
CHARCOAL = RGBColor(*hex_to_rgb('#2D3748'))
TEAL = RGBColor(*hex_to_rgb('#14B8A6'))
TEXT = RGBColor(*hex_to_rgb('#1A202C'))
TEXT_LIGHT = RGBColor(*hex_to_rgb('#4A5568'))
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(*hex_to_rgb('#F8F9FA'))
BORDER_GRAY = RGBColor(*hex_to_rgb('#E5E7EB'))

def create_slide1_default():
    """Create Slide 1 DEFAULT - Balance Scale Imbalance"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Create blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Margins: 80px = ~0.83 inches
    MARGIN = Inches(0.83)
    CONTENT_WIDTH = Inches(16) - (2 * MARGIN)
    CONTENT_HEIGHT = Inches(9) - (2 * MARGIN)

    # Headline
    headline = slide.shapes.add_textbox(
        MARGIN, MARGIN,
        CONTENT_WIDTH, Inches(1)
    )
    tf = headline.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Command's Market Presence is Strong, But Creator Discussions Reveal a Tension"
    p.font.name = 'Inter'
    p.font.size = Pt(33)
    p.font.bold = True
    p.font.color.rgb = TEXT

    # Left Column (800px = ~8.33 inches)
    LEFT_COL_WIDTH = Inches(8.33)
    LEFT_COL_X = MARGIN
    LEFT_COL_Y = MARGIN + Inches(1.2)

    # Pattern Card
    pattern_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_COL_X, LEFT_COL_Y,
        LEFT_COL_WIDTH, Inches(1.5)
    )
    pattern_card.fill.solid()
    pattern_card.fill.fore_color.rgb = LIGHT_BG
    pattern_card.line.color.rgb = BORDER_GRAY
    pattern_card.line.width = Pt(2)

    # Pattern text
    tf = pattern_card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(16)
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)

    # Title
    p = tf.paragraphs[0]
    p.text = "THE PATTERN"
    p.font.name = 'Inter'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL

    # Body
    p2 = tf.add_paragraph()
    p2.text = "Video creators frequently highlight Command's surface damage prevention (24 mentions), positioning it as ideal for rental properties and temporary installations. However, surface damage complaints (24 severity) equally appear across content, creating cognitive dissonance between marketing promise and user experience."
    p2.font.name = 'Inter'
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT
    p2.line_spacing = Pt(19)

    # Insight Card
    insight_y = LEFT_COL_Y + Inches(1.7)
    insight_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_COL_X, insight_y,
        LEFT_COL_WIDTH, Inches(1.5)
    )
    insight_card.fill.solid()
    insight_card.fill.fore_color.rgb = LIGHT_BG
    insight_card.line.color.rgb = BORDER_GRAY
    insight_card.line.width = Pt(2)

    tf = insight_card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(16)
    tf.margin_left = Pt(20)
    tf.margin_right = Pt(20)

    p = tf.paragraphs[0]
    p.text = "INSIGHT"
    p.font.name = 'Inter'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL

    p2 = tf.add_paragraph()
    p2.text = "The tension isn't about product failure—it's about expectation management. Command's 'damage-free' positioning sets absolute expectations, so any surface damage feels like brand promise betrayal rather than acceptable risk."
    p2.font.name = 'Inter'
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT
    p2.line_spacing = Pt(19)

    # Hard Truth Card (dark gradient)
    hardtruth_y = insight_y + Inches(1.7)
    hardtruth_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        LEFT_COL_X, hardtruth_y,
        LEFT_COL_WIDTH, Inches(1.8)
    )
    hardtruth_card.fill.solid()
    hardtruth_card.fill.fore_color.rgb = CHARCOAL
    hardtruth_card.line.color.rgb = TEAL
    hardtruth_card.line.width = Pt(6)

    tf = hardtruth_card.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(18)
    tf.margin_left = Pt(22)
    tf.margin_right = Pt(22)

    p = tf.paragraphs[0]
    p.text = "HARD TRUTH"
    p.font.name = 'Inter'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEAL

    p2 = tf.add_paragraph()
    p2.text = "Equal mention of 'surface damage' and 'damage-free' creates brand credibility crisis."
    p2.font.name = 'Inter'
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.line_spacing = Pt(21)

    p3 = tf.add_paragraph()
    p3.text = "Strategic choice: Either fix product formulation (costly, R&D-intensive) or transparently communicate surface limitations on packaging (honesty over false promise). Current positioning gap drives users to traditional nails/screws where damage is expected and controlled."
    p3.font.name = 'Inter'
    p3.font.size = Pt(11)
    p3.font.color.rgb = WHITE
    p3.line_spacing = Pt(18)

    # Footer
    footer = slide.shapes.add_textbox(
        MARGIN, Inches(9) - MARGIN - Inches(0.5),
        CONTENT_WIDTH, Inches(0.4)
    )
    footer.line.color.rgb = TEAL
    footer.line.width = Pt(2)

    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Source: Phase 2 Analysis, Command brand perception across 62 creator videos"
    p.font.name = 'Inter'
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT

    return prs

def main():
    print("🎯 Converting Slide 1 DEFAULT to PowerPoint...")
    prs = create_slide1_default()
    output_path = "SLIDE1_DEFAULT.pptx"
    prs.save(output_path)
    print(f"✅ Saved: {output_path}")
    print(f"📊 Slide dimensions: 16:9 (1920×1080)")
    print(f"🎨 Colors: Charcoal (#2D3748), Teal (#14B8A6)")
    print(f"📝 Typography: Inter font preserved")

if __name__ == "__main__":
    main()
