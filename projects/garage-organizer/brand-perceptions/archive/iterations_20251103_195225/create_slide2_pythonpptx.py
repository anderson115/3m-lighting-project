"""Recreate Slide 2 (Default) as editable PowerPoint using python-pptx."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Unit helpers
PX_TO_IN = 1 / 120  # HTML canvas: 1920px ⇒ 16in, so 120px per inch


def px(value: float) -> float:
    """Convert pixel values from the HTML canvas to inches."""
    return Inches(value * PX_TO_IN)


# Color palette
CHARCOAL = RGBColor(45, 55, 72)        # #2D3748 primary
ACCENT = RGBColor(20, 184, 166)        # #14B8A6 teal
TEXT = RGBColor(26, 32, 44)            # #1A202C
TEXT_LIGHT = RGBColor(100, 116, 139)   # #64748B
GRAY_BG = RGBColor(229, 231, 235)      # #E5E7EB
MEDIUM_GRAY = RGBColor(148, 163, 184)  # #94A3B8
DARK_GRAY = RGBColor(71, 85, 105)      # #475569
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(245, 158, 11)        # #F59E0B
PALE_ORANGE = RGBColor(253, 230, 138)  # #FDE68A
ORANGE_TEXT = RGBColor(146, 64, 14)    # #92400E

# Layout helpers based on HTML pixel grid
SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)
MARGIN = px(80)
LEFT_COL_WIDTH = px(800)
RIGHT_COL_WIDTH = px(900)
GAP = px(60)
CONTENT_WIDTH = px(1760)


def add_text_box(slide, left, top, width, height, text, *, font_size, bold=False,
                 color=TEXT, line_spacing=None, uppercase=False, align=PP_ALIGN.LEFT,
                 spacing_after=0, spacing_before=0):
    """Create a textbox with single paragraph of formatted text."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.upper() if uppercase else text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Inter"
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    p.space_after = Pt(spacing_after)
    p.space_before = Pt(spacing_before)
    return shape


def add_body_text(slide, left, top, width, text):
    shape = slide.shapes.add_textbox(left, top, width, Inches(2))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Inter"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT
    p.line_spacing = 1.5
    return shape


def add_section(slide, top, header, body):
    add_text_box(
        slide,
        MARGIN,
        top,
        LEFT_COL_WIDTH,
        Inches(0.3),
        header,
        font_size=11,
        bold=True,
        color=CHARCOAL,
        uppercase=True,
        spacing_after=2,
    )
    add_body_text(slide, MARGIN, top + Inches(0.25), LEFT_COL_WIDTH, body)


def add_hard_truth(slide, top, title, statement, explanation):
    height = Inches(2.05)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN,
        top,
        LEFT_COL_WIDTH,
        height,
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(17, 24, 39)  # deep navy backdrop
    panel.line.color.rgb = ACCENT
    panel.line.width = Pt(5)
    panel.shadow.inherit = False

    # Diamond icon using small rotated square
    diamond = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        MARGIN + Inches(0.18),
        top + Inches(0.3),
        Inches(0.25),
        Inches(0.25),
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = ACCENT
    diamond.line.color.rgb = ACCENT

    text_box = slide.shapes.add_textbox(
        MARGIN + Inches(0.55),
        top + Inches(0.25),
        LEFT_COL_WIDTH - Inches(0.8),
        height - Inches(0.5),
    )
    tf = text_box.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.name = "Inter"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = statement
    p.font.name = "Inter"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(6)
    p.line_spacing = 1.4

    p = tf.add_paragraph()
    p.text = explanation
    p.font.name = "Inter"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.space_before = Pt(6)
    p.line_spacing = 1.4


def add_provocation(slide, top, question, items):
    # Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        MARGIN,
        top,
        Inches(0.4),
        Inches(0.4),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.color.rgb = CHARCOAL
    circle.line.width = Pt(2)

    # Question text
    add_text_box(
        slide,
        MARGIN + Inches(0.45),
        top - Inches(0.05),
        LEFT_COL_WIDTH - Inches(0.45),
        Inches(0.6),
        question,
        font_size=13,
        bold=True,
        color=TEXT,
    )

    # Subheader
    add_text_box(
        slide,
        MARGIN,
        top + Inches(0.55),
        LEFT_COL_WIDTH,
        Inches(0.2),
        "Consider:",
        font_size=11,
        bold=True,
        color=CHARCOAL,
    )

    # Bullet list
    list_box = slide.shapes.add_textbox(
        MARGIN,
        top + Inches(0.75),
        LEFT_COL_WIDTH,
        Inches(1.2),
    )
    tf = list_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.add_paragraph() if idx else tf.paragraphs[0]
        p.text = f"▸ {item}"
        p.font.name = "Inter"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT
        p.line_spacing = 1.4
        p.space_before = Pt(2)


def add_performance_item(slide, top, label, status, value, width_pct):
    col_left = MARGIN + LEFT_COL_WIDTH + GAP
    label_height = px(26)
    status_width = px(220)
    label_width = RIGHT_COL_WIDTH - status_width - px(12)

    # Label (left aligned)
    label_box = slide.shapes.add_textbox(
        col_left,
        top,
        label_width,
        label_height,
    )
    ltf = label_box.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.name = "Inter"
    lp.font.size = Pt(11)
    lp.font.bold = True
    lp.font.color.rgb = TEXT

    # Status (right aligned)
    status_box = slide.shapes.add_textbox(
        col_left + label_width + px(12),
        top,
        status_width,
        label_height,
    )
    stf = status_box.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = status
    sp.font.name = "Inter"
    sp.font.size = Pt(11)
    sp.font.bold = True
    sp.font.color.rgb = ACCENT if float(width_pct.strip('%')) > 60 else DARK_GRAY
    sp.alignment = PP_ALIGN.RIGHT

    # Bar container beneath labels
    bar_top = top + label_height + px(6)
    bar_height = px(32)
    container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        col_left,
        bar_top,
        RIGHT_COL_WIDTH,
        bar_height,
    )
    container.fill.solid()
    container.fill.fore_color.rgb = GRAY_BG
    container.line.color.rgb = CHARCOAL
    container.line.width = Pt(2)

    bar_pct = max(min(float(width_pct.strip('%')), 100.0), 0.0) / 100.0
    inner_width = max(RIGHT_COL_WIDTH * bar_pct - px(16), px(60))
    bar_color = ACCENT if bar_pct >= 0.6 else MEDIUM_GRAY
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        col_left + px(4),
        bar_top + px(4),
        inner_width,
        bar_height - px(8),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.color.rgb = bar_color

    value_box = slide.shapes.add_textbox(
        bar.left + bar.width - px(120),
        bar_top + px(2),
        px(110),
        bar_height - px(4),
    )
    vtf = value_box.text_frame
    vtf.clear()
    vp = vtf.paragraphs[0]
    vp.text = value
    vp.font.name = "Inter"
    vp.font.size = Pt(12)
    vp.font.bold = True
    vp.font.color.rgb = WHITE
    vp.alignment = PP_ALIGN.RIGHT

    return bar_top + bar_height


def add_sentiment_cards(slide, top, sentiments):
    card_width = (RIGHT_COL_WIDTH - Inches(0.3)) / 2
    for idx, (value, label, is_negative) in enumerate(sentiments):
        left = MARGIN + LEFT_COL_WIDTH + GAP + idx * (card_width + Inches(0.3))
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left,
            top,
            card_width,
            Inches(1.6),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        border_color = CHARCOAL if is_negative else ACCENT
        card.line.color.rgb = border_color
        card.line.width = Pt(3)

        tf = card.text_frame
        tf.clear()
        tf.vertical_anchor = 1
        value_p = tf.paragraphs[0]
        value_p.text = value
        value_p.font.name = "Inter"
        value_p.font.size = Pt(30)
        value_p.font.bold = True
        value_p.font.color.rgb = border_color
        value_p.alignment = PP_ALIGN.CENTER

        label_p = tf.add_paragraph()
        label_p.text = label.upper()
        label_p.font.name = "Inter"
        label_p.font.size = Pt(10)
        label_p.font.bold = True
        label_p.font.color.rgb = TEXT
        label_p.alignment = PP_ALIGN.CENTER
        label_p.space_before = Pt(8)


def add_temperature_card(slide, top, stat, description):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN + LEFT_COL_WIDTH + GAP,
        top,
        RIGHT_COL_WIDTH,
        Inches(1.4),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = PALE_ORANGE
    card.line.color.rgb = ORANGE
    card.line.width = Pt(4)

    header = slide.shapes.add_textbox(
        card.left + Inches(0.3),
        top + Inches(0.2),
        RIGHT_COL_WIDTH - Inches(0.6),
        Inches(0.25),
    )
    htf = header.text_frame
    htf.text = "Temperature Discussion Frequency".upper()
    hp = htf.paragraphs[0]
    hp.font.size = Pt(10)
    hp.font.bold = True
    hp.font.color.rgb = ORANGE_TEXT
    hp.font.name = "Inter"

    stat_box = slide.shapes.add_textbox(
        card.left + Inches(0.3),
        top + Inches(0.5),
        RIGHT_COL_WIDTH - Inches(0.6),
        Inches(0.4),
    )
    stf = stat_box.text_frame
    stf.text = stat
    sp = stf.paragraphs[0]
    sp.font.size = Pt(28)
    sp.font.bold = True
    sp.font.color.rgb = ORANGE
    sp.font.name = "Inter"
    sp.alignment = PP_ALIGN.LEFT

    desc_shape = slide.shapes.add_textbox(
        card.left + Inches(0.3),
        top + Inches(0.95),
        RIGHT_COL_WIDTH - Inches(0.6),
        Inches(0.35),
    )
    tf = desc_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = description
    p.font.name = "Inter"
    p.font.size = Pt(10)
    p.font.color.rgb = ORANGE_TEXT
    p.line_spacing = 1.3

    return card.top + card.height


def add_conversation(slide, top, items):
    header = slide.shapes.add_textbox(
        MARGIN,
        top,
        CONTENT_WIDTH,
        Inches(0.3),
    )
    htf = header.text_frame
    htf.text = "Conversation Starters".upper()
    hp = htf.paragraphs[0]
    hp.font.size = Pt(11)
    hp.font.bold = True
    hp.font.color.rgb = CHARCOAL
    hp.font.name = "Inter"

    list_box = slide.shapes.add_textbox(
        MARGIN,
        top + Inches(0.25),
        CONTENT_WIDTH,
        Inches(0.9),
    )
    tf = list_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.add_paragraph() if idx else tf.paragraphs[0]
        p.text = f"→ {item}"
        p.font.name = "Inter"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT
        p.line_spacing = 1.4
        p.space_before = Pt(2)


def build_slide():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Headline
    headline_text = (
        "Scotch Demonstrates Strong Creator Advocacy—Understanding Where It Shines "
        "and Where It's Challenged"
    )
    headline_box = add_text_box(
        slide,
        MARGIN,
        MARGIN,
        CONTENT_WIDTH,
        Inches(1),
        headline_text,
        font_size=30,
        bold=True,
        color=TEXT,
    )
    headline_box.text_frame.paragraphs[0].line_spacing = 1.2

    # Sections
    top = MARGIN + Inches(1.05)
    add_section(
        slide,
        top,
        "✱ The Pattern",
        (
            "Content creators who feature Scotch products tend to offer favorable "
            "recommendations. Across review and comparison content, Scotch receives "
            "notably more positive assessments than negative, with creators frequently "
            "emphasizing reliability and hold strength."
        ),
    )

    top += Inches(1.4)
    add_section(
        slide,
        top,
        "✱ Insight",
        (
            "Scotch's performance perception is strongest in indoor, permanent mounting "
            "applications. Temperature considerations shape use case recommendations for "
            "outdoor/vehicle contexts—creators understand optimal use through experience."
        ),
    )

    top += Inches(1.5)
    add_hard_truth(
        slide,
        top,
        "Hard Truth",
        "Temperature boundary discussions signal market opportunity cost—outdoor/vehicle segments constrained.",
        (
            "Strategic choice: Invest in all-temperature formulation to expand TAM, or "
            "own indoor category dominance and accept boundary? Current positioning leaves "
            "outdoor segment vulnerable to competitors."
        ),
    )

    top += Inches(2.3)
    add_provocation(
        slide,
        top,
        "Is temperature sensitivity fixable or should we communicate boundaries transparently?",
        [
            "Product innovation: All-temperature line for outdoor segment",
            "Transparent positioning: \"Optimized for indoor use, 40-100°F\"",
            "Accept boundary and focus on indoor permanent mounting dominance",
        ],
    )

    # Right column: Performance spectrum
    perf_top = MARGIN + Inches(1.1)
    add_text_box(
        slide,
        MARGIN + LEFT_COL_WIDTH + GAP,
        perf_top,
        RIGHT_COL_WIDTH,
        px(28),
        "Performance Perception by Context",
        font_size=11,
        bold=True,
        color=CHARCOAL,
        uppercase=True,
    )

    y = perf_top + px(36)
    y = add_performance_item(
        slide,
        y,
        "Indoor Permanent Mounting",
        "High Confidence",
        "92%",
        "92%",
    ) + px(24)
    y = add_performance_item(
        slide,
        y,
        "DIY Project Applications",
        "Strong",
        "88%",
        "88%",
    ) + px(24)
    y = add_performance_item(
        slide,
        y,
        "Outdoor/Extreme Temperature",
        "Consideration",
        "45%",
        "45%",
    )

    # Sentiment cards
    sentiment_top = y + px(54)
    add_text_box(
        slide,
        MARGIN + LEFT_COL_WIDTH + GAP,
        sentiment_top,
        RIGHT_COL_WIDTH,
        px(28),
        "Creator Sentiment Analysis",
        font_size=11,
        bold=True,
        color=CHARCOAL,
        uppercase=True,
    )
    add_sentiment_cards(
        slide,
        sentiment_top + px(36),
        [
            ("High", "Positive Assessments", False),
            ("Low", "Negative Feedback", True),
        ],
    )

    # Temperature indicator beneath sentiment cards
    temp_bottom = add_temperature_card(
        slide,
        sentiment_top + px(36) + Inches(1.6) + px(36),
        "1 in 6",
        "Temperature performance mentioned in ~17% of Scotch content, primarily outdoor/vehicle contexts",
    )

    # Conversation starters (full width)
    conversation_block_height = Inches(0.3) + Inches(0.9)
    footer_height = Inches(0.4)
    max_conversation_top = SLIDE_HEIGHT - conversation_block_height - footer_height - px(24)
    conversation_top = min(temp_bottom + px(36), max_conversation_top)
    add_conversation(
        slide,
        conversation_top,
        [
            "Could an 'All-Temperature' product line address outdoor segments without cannibalizing indoor business?",
            "Should we make temperature range explicit on packaging to set proper expectations?",
            "DIY and maker communities are growing—how do we deepen presence here as \"professional choice\"?",
        ],
    )

    # Footer
    footer_top = SLIDE_HEIGHT - footer_height - px(12)
    footer = slide.shapes.add_textbox(
        MARGIN,
        footer_top,
        SLIDE_WIDTH - 2 * MARGIN,
        footer_height,
    )
    footer.line.color.rgb = ACCENT
    footer.line.width = Pt(2)
    tf = footer.text_frame
    tf.text = "Source: Phase 2 Analysis, Scotch performance perception across 58 creator videos"
    p = tf.paragraphs[0]
    p.font.name = "Inter"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT

    output = Path(__file__).with_name("SLIDE2_pythonpptx.pptx")
    prs.save(output)
    print(f"Saved editable slide to {output}")


if __name__ == "__main__":
    build_slide()
# Slide geometry
SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)
