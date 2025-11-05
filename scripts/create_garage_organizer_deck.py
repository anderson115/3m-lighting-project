"""
Create 3M Garage Organization Category Intelligence Deck
Complete client-ready PowerPoint presentation with content from research reports
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === COLOR PALETTE (offbrain Design System) ===
CHARCOAL = RGBColor(45, 55, 72)        # #2D3748 primary
ACCENT = RGBColor(20, 184, 166)        # #14B8A6 teal
TEXT = RGBColor(26, 32, 44)            # #1A202C dark gray
TEXT_LIGHT = RGBColor(74, 85, 104)     # #4A5568 light gray
BG_PANEL = RGBColor(249, 250, 251)     # #F9FAFB
WHITE = RGBColor(255, 255, 255)
DEEP_NAVY = RGBColor(17, 24, 39)       # #111827
AMBER = RGBColor(245, 158, 11)         # #F59E0B
AMBER_LIGHT = RGBColor(254, 243, 199)  # #FEF3C7
AMBER_TEXT = RGBColor(146, 64, 14)     # #92400E

# === LAYOUT SPECIFICATIONS ===
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)
MARGIN = Inches(0.83)  # 80px at 96 DPI
LEFT_COL_WIDTH = Inches(4.17)  # 800px
RIGHT_COL_WIDTH = Inches(4.69)  # 900px
GAP = Inches(0.31)  # 60px

def create_presentation():
    """Create blank presentation with 16:9 aspect ratio."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs

def add_text_box(slide, left, top, width, height, text, **kwargs):
    """Add formatted text box to slide."""
    font_size = kwargs.get('font_size', 12)
    bold = kwargs.get('bold', False)
    color = kwargs.get('color', TEXT)
    align = kwargs.get('align', PP_ALIGN.LEFT)
    line_spacing = kwargs.get('line_spacing', 1.0)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Inter"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = line_spacing
    return shape

def add_headline(slide, text):
    """Add slide headline."""
    return add_text_box(
        slide, MARGIN, MARGIN, Inches(9.17), Inches(0.6),
        text, font_size=33, bold=True, color=TEXT, line_spacing=1.15
    )

def add_section_header(slide, left, top, text):
    """Add uppercase section header."""
    tb = add_text_box(
        slide, left, top, Inches(4), Inches(0.25),
        text.upper(), font_size=12, bold=True, color=CHARCOAL
    )
    tb.text_frame.paragraphs[0].font.letter_spacing = Pt(1)
    return tb

def add_body_text(slide, left, top, width, text):
    """Add body text paragraph."""
    return add_text_box(
        slide, left, top, width, Inches(1.5),
        text, font_size=12, color=TEXT, line_spacing=1.58
    )

def add_footer(slide, source_text):
    """Add footer with source citation."""
    add_text_box(
        slide, MARGIN, Inches(5.2), Inches(9), Inches(0.2),
        f"Source: {source_text}",
        font_size=10, color=TEXT_LIGHT
    )

def add_hard_truth_panel(slide, left, top, width, statement, explanation):
    """Add Hard Truth callout panel with diamond icon."""
    height = Inches(1.2)

    # Background panel with gradient effect
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = DEEP_NAVY
    panel.line.color.rgb = ACCENT
    panel.line.width = Pt(5)
    panel.shadow.inherit = False

    # Diamond icon
    diamond = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        left + Inches(0.15), top + Inches(0.18),
        Inches(0.19), Inches(0.19)
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = ACCENT
    diamond.line.fill.background()

    # Header
    add_text_box(
        slide, left + Inches(0.4), top + Inches(0.15), width - Inches(0.5), Inches(0.2),
        "HARD TRUTH", font_size=10, bold=True, color=ACCENT
    )

    # Statement
    add_text_box(
        slide, left + Inches(0.4), top + Inches(0.38), width - Inches(0.5), Inches(0.35),
        statement, font_size=14, bold=True, color=WHITE, line_spacing=1.4
    )

    # Explanation
    add_text_box(
        slide, left + Inches(0.4), top + Inches(0.75), width - Inches(0.5), Inches(0.35),
        explanation, font_size=10, color=WHITE, line_spacing=1.6
    )

def add_data_card(slide, left, top, width, height, value, label, accent_color=ACCENT):
    """Add data visualization card with large number."""
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_PANEL
    card.line.fill.background()

    # Large value
    add_text_box(
        slide, left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), Inches(0.4),
        value, font_size=36, bold=True, color=accent_color, align=PP_ALIGN.CENTER
    )

    # Label
    add_text_box(
        slide, left + Inches(0.1), top + height - Inches(0.35), width - Inches(0.2), Inches(0.3),
        label, font_size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER
    )

def create_slide_1_executive_summary(prs):
    """Slide 1: Executive Summary with dataset stats and key insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    add_headline(slide, "3M Garage Organization Market Intelligence: Executive Summary")

    # Dataset stats (4 cards)
    stats = [
        ("9,555", "Products Analyzed"),
        ("5", "Major Retailers"),
        ("571+", "Consumer Videos"),
        ("$518K", "Monthly Revenue\n(Top 20 SKUs)")
    ]

    card_width = Inches(2.0)
    card_height = Inches(0.9)
    start_left = MARGIN
    top = MARGIN + Inches(0.8)

    for i, (value, label) in enumerate(stats):
        add_data_card(
            slide,
            start_left + i * (card_width + Inches(0.15)),
            top,
            card_width,
            card_height,
            value,
            label
        )

    # Key Insights section
    insights_top = top + Inches(1.2)
    add_section_header(slide, MARGIN, insights_top, "KEY INSIGHTS")

    insights = [
        ("Quality Crisis", "90% of best-selling products generate negative quality sentiment despite $518K monthly revenue"),
        ("Premium Gap", "Less than 3% of market in $40-80 premium segment presents untapped opportunity"),
        ("Technology Advantage", "VHB™ adhesive addresses #1 consumer pain point: damage-free mounting"),
        ("Purchase Evolution", "Customers shift from price-focus (34%) to quality-focus (41%) after initial purchase")
    ]

    insight_top = insights_top + Inches(0.3)
    for i, (title, desc) in enumerate(insights):
        y = insight_top + i * Inches(0.5)
        add_text_box(
            slide, MARGIN, y, Inches(9), Inches(0.15),
            f"{i+1}. {title}:", font_size=11, bold=True, color=CHARCOAL
        )
        add_text_box(
            slide, MARGIN + Inches(0.2), y + Inches(0.18), Inches(8.8), Inches(0.3),
            desc, font_size=10, color=TEXT, line_spacing=1.4
        )

    add_footer(slide, "3M Category Intelligence Analysis | October 2025")
    return slide

def create_slide_2_market_opportunity(prs):
    """Slide 2: Market Opportunity - Quality Gap Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_headline(slide, "Massive Quality Gap Despite Strong Sales Performance")

    # Left column - Narrative
    left_top = MARGIN + Inches(0.8)

    add_section_header(slide, MARGIN, left_top, "THE PATTERN")
    add_body_text(
        slide, MARGIN, left_top + Inches(0.3), LEFT_COL_WIDTH,
        "Market demonstrates concentrated demand with top 20 SKUs generating 28,770 units/month, yet 90% of best-sellers produce negative quality sentiment."
    )

    add_section_header(slide, MARGIN, left_top + Inches(1.0), "THE INSIGHT")
    add_body_text(
        slide, MARGIN, left_top + Inches(1.3), LEFT_COL_WIDTH,
        "Consumers accept poor quality due to lack of premium alternatives, creating opportunity for differentiated positioning."
    )

    # Hard Truth Panel
    add_hard_truth_panel(
        slide, MARGIN, left_top + Inches(2.2), LEFT_COL_WIDTH,
        "Current market accepts systemic failure as normal",
        "Products routinely fail at 30% of rated capacity. Catastrophic mounting failures cause property damage. Premium segment is virtually non-existent."
    )

    # Right column - Data visualization
    right_left = MARGIN + LEFT_COL_WIDTH + GAP

    # Top 20 Performance Card
    add_data_card(
        slide, right_left, left_top, RIGHT_COL_WIDTH * 0.48, Inches(1.0),
        "28,770", "Units/Month\n(Top 20 SKUs)"
    )

    # Revenue Card
    add_data_card(
        slide, right_left + RIGHT_COL_WIDTH * 0.52, left_top, RIGHT_COL_WIDTH * 0.48, Inches(1.0),
        "$518K", "Monthly Revenue\n(Top 20 SKUs)"
    )

    # Quality Sentiment Bar
    sentiment_top = left_top + Inches(1.3)
    add_text_box(
        slide, right_left, sentiment_top, RIGHT_COL_WIDTH, Inches(0.2),
        "Quality Sentiment Distribution", font_size=11, bold=True, color=CHARCOAL
    )

    # Negative bar
    neg_width = RIGHT_COL_WIDTH * 0.90
    bar_height = Inches(0.35)
    neg_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        right_left, sentiment_top + Inches(0.3), neg_width, bar_height
    )
    neg_bar.fill.solid()
    neg_bar.fill.fore_color.rgb = RGBColor(239, 68, 68)  # Red
    neg_bar.line.fill.background()

    add_text_box(
        slide, right_left + Inches(0.1), sentiment_top + Inches(0.38), neg_width, Inches(0.2),
        "90% Negative Sentiment", font_size=12, bold=True, color=WHITE
    )

    # Positive bar
    pos_width = RIGHT_COL_WIDTH * 0.10
    pos_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        right_left + neg_width, sentiment_top + Inches(0.3), pos_width, bar_height
    )
    pos_bar.fill.solid()
    pos_bar.fill.fore_color.rgb = RGBColor(34, 197, 94)  # Green
    pos_bar.line.fill.background()

    # Failure modes table
    table_top = sentiment_top + Inches(0.9)
    add_text_box(
        slide, right_left, table_top, RIGHT_COL_WIDTH, Inches(0.2),
        "Primary Failure Categories", font_size=11, bold=True, color=CHARCOAL
    )

    failures = [
        ("Weight Capacity Failures", "67%"),
        ("Mounting System Failures", "41%"),
        ("Durability Issues", "38%")
    ]

    for i, (failure, pct) in enumerate(failures):
        y = table_top + Inches(0.3) + i * Inches(0.28)
        add_text_box(
            slide, right_left, y, RIGHT_COL_WIDTH * 0.7, Inches(0.2),
            failure, font_size=10, color=TEXT
        )
        add_text_box(
            slide, right_left + RIGHT_COL_WIDTH * 0.75, y, RIGHT_COL_WIDTH * 0.25, Inches(0.2),
            pct, font_size=12, bold=True, color=RGBColor(239, 68, 68), align=PP_ALIGN.RIGHT
        )

    add_footer(slide, "Analysis of 2,847 negative reviews across 5 retail channels")
    return slide

def create_slide_3_competitive_positioning(prs):
    """Slide 3: Competitive Positioning Matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_headline(slide, "Strategic Premium Gap in Competitive Landscape")

    # Left column
    left_top = MARGIN + Inches(0.8)

    add_section_header(slide, MARGIN, left_top, "MARKET STRUCTURE")
    add_body_text(
        slide, MARGIN, left_top + Inches(0.3), LEFT_COL_WIDTH,
        "85% of market concentrated in commodity segment ($5-20), with premium tier ($40-80) representing less than 3% of products but demonstrating higher satisfaction rates."
    )

    add_hard_truth_panel(
        slide, MARGIN, left_top + Inches(1.2), LEFT_COL_WIDTH,
        "$40-80 premium gap is virtually uncontested",
        "Current market structure creates race-to-bottom in commodity tier while professional tier ($80+) remains fragmented. White space opportunity for premium performance positioning."
    )

    # Right column - Positioning table
    right_left = MARGIN + LEFT_COL_WIDTH + GAP

    add_text_box(
        slide, right_left, left_top, RIGHT_COL_WIDTH, Inches(0.25),
        "MARKET SEGMENTATION", font_size=12, bold=True, color=CHARCOAL
    )

    # Table headers
    headers = ["Segment", "Price Range", "Share", "Key Players"]
    header_y = left_top + Inches(0.35)
    col_widths = [Inches(1.3), Inches(1.0), Inches(0.8), Inches(1.5)]

    x_pos = right_left
    for i, header in enumerate(headers):
        add_text_box(
            slide, x_pos, header_y, col_widths[i], Inches(0.2),
            header, font_size=9, bold=True, color=CHARCOAL
        )
        x_pos += col_widths[i]

    # Table rows
    rows = [
        ("Commodity", "$5-20", "85%", "Rubbermaid, Everbilt"),
        ("Value", "$20-40", "12%", "Gladiator, Husky"),
        ("Premium Gap", "$40-80", "<3%", "3M OPPORTUNITY"),
        ("Professional", "$80+", "<1%", "StoreWall, Monkey Bar")
    ]

    row_y = header_y + Inches(0.3)
    for row_data in rows:
        x_pos = right_left
        is_opportunity = "OPPORTUNITY" in row_data[3]
        bg_color = AMBER_LIGHT if is_opportunity else BG_PANEL
        text_color = CHARCOAL if is_opportunity else TEXT

        # Background for opportunity row
        if is_opportunity:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                right_left - Inches(0.05), row_y - Inches(0.05),
                sum(col_widths) + Inches(0.1), Inches(0.28)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = bg_color
            bg.line.fill.background()

        for i, value in enumerate(row_data):
            add_text_box(
                slide, x_pos, row_y, col_widths[i], Inches(0.2),
                value, font_size=9 if not is_opportunity else 10,
                bold=is_opportunity, color=text_color
            )
            x_pos += col_widths[i]
        row_y += Inches(0.3)

    add_footer(slide, "Market share analysis from 9,555 products across 5 retailers")
    return slide

def create_slide_4_technology_advantage(prs):
    """Slide 4: 3M Technology Advantage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_headline(slide, "3M Technology Directly Addresses Top Consumer Pain Points")

    # Left column
    left_top = MARGIN + Inches(0.8)

    add_section_header(slide, MARGIN, left_top, "VHB™ ADHESIVE TECHNOLOGY")
    add_body_text(
        slide, MARGIN, left_top + Inches(0.3), LEFT_COL_WIDTH,
        "Addresses #1 consumer pain point: mounting without wall damage. 10x stronger than mechanical fasteners, temperature stable -40°F to 200°F, proven in architectural applications supporting 1000+ lbs."
    )

    add_section_header(slide, MARGIN, left_top + Inches(1.1), "ADVANCED MATERIALS SCIENCE")
    add_body_text(
        slide, MARGIN, left_top + Inches(1.4), LEFT_COL_WIDTH,
        "Powder coating expertise eliminates rust issues (38% of current complaints). Composite materials for weight optimization. Surface treatments for enhanced durability."
    )

    add_hard_truth_panel(
        slide, MARGIN, left_top + Inches(2.3), LEFT_COL_WIDTH,
        "Technology advantage creates sustainable moat",
        "Competitors lack adhesive expertise and materials science capabilities. 3M's existing retail relationships accelerate distribution. Scale manufacturing maintains quality at volume."
    )

    # Right column - Performance bars
    right_left = MARGIN + LEFT_COL_WIDTH + GAP

    add_text_box(
        slide, right_left, left_top, RIGHT_COL_WIDTH, Inches(0.25),
        "COMPETITIVE ADVANTAGE METRICS", font_size=12, bold=True, color=CHARCOAL
    )

    # Performance bars
    metrics = [
        ("VHB™ Adhesion Strength", 95, "10x vs. mechanical fasteners"),
        ("Temperature Stability", 92, "-40°F to 200°F range"),
        ("Durability/Anti-Rust", 88, "Powder coating expertise")
    ]

    bar_top = left_top + Inches(0.4)
    for i, (label, score, detail) in enumerate(metrics):
        y = bar_top + i * Inches(0.7)

        # Label
        add_text_box(
            slide, right_left, y, RIGHT_COL_WIDTH, Inches(0.2),
            label, font_size=11, bold=True, color=TEXT
        )

        # Background bar (gray)
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            right_left, y + Inches(0.25), RIGHT_COL_WIDTH, Inches(0.25)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(229, 231, 235)
        bg_bar.line.fill.background()

        # Score bar (teal)
        score_width = RIGHT_COL_WIDTH * (score / 100)
        score_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            right_left, y + Inches(0.25), score_width, Inches(0.25)
        )
        score_bar.fill.solid()
        score_bar.fill.fore_color.rgb = ACCENT
        score_bar.line.fill.background()

        # Score text
        add_text_box(
            slide, right_left + Inches(0.1), y + Inches(0.28), score_width, Inches(0.2),
            f"{score}%", font_size=11, bold=True, color=WHITE
        )

        # Detail text
        add_text_box(
            slide, right_left, y + Inches(0.52), RIGHT_COL_WIDTH, Inches(0.15),
            detail, font_size=9, color=TEXT_LIGHT
        )

    add_footer(slide, "Technology assessment based on 3M materials science portfolio")
    return slide

def create_slide_5_product_roadmap(prs):
    """Slide 5: Product Development Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_headline(slide, "Three-Phase Product Development Roadmap")

    left_top = MARGIN + Inches(0.8)

    # Phase cards
    phases = [
        {
            "phase": "PHASE 1: HERO PRODUCT",
            "timeframe": "Months 0-6",
            "product": "VHB™ Heavy-Duty Hook System",
            "details": "3 SKUs (25/50/100 lb capacity)\nRetail: $49-89 | Margin: 65%\nHome Depot exclusive launch",
            "color": ACCENT
        },
        {
            "phase": "PHASE 2: CATEGORY EXPANSION",
            "timeframe": "Months 7-12",
            "product": "Modular Systems",
            "details": "Rail System ($129-249)\nOverhead Storage ($199-299)\nSpecialty Solutions (Bike, Sports, Tools)",
            "color": CHARCOAL
        },
        {
            "phase": "PHASE 3: ECOSYSTEM",
            "timeframe": "Year 2",
            "product": "Smart Features + Services",
            "details": "Smart weight sensors with app\nSubscription consumables\nProfessional installation\nB2B commercial segment",
            "color": AMBER
        }
    ]

    card_height = Inches(1.3)
    for i, phase_data in enumerate(phases):
        y = left_top + i * (card_height + Inches(0.15))

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN, y, Inches(9), card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = phase_data["color"]
        card.line.width = Pt(3)

        # Phase header
        add_text_box(
            slide, MARGIN + Inches(0.2), y + Inches(0.15), Inches(4), Inches(0.2),
            phase_data["phase"], font_size=12, bold=True, color=phase_data["color"]
        )

        # Timeframe
        add_text_box(
            slide, MARGIN + Inches(0.2), y + Inches(0.38), Inches(2), Inches(0.15),
            phase_data["timeframe"], font_size=10, color=TEXT_LIGHT
        )

        # Product name
        add_text_box(
            slide, MARGIN + Inches(0.2), y + Inches(0.6), Inches(4), Inches(0.2),
            phase_data["product"], font_size=14, bold=True, color=TEXT
        )

        # Details (right side)
        add_text_box(
            slide, MARGIN + Inches(4.5), y + Inches(0.38), Inches(4.3), Inches(0.8),
            phase_data["details"], font_size=10, color=TEXT, line_spacing=1.5
        )

    add_footer(slide, "Product development timeline from 03_PRODUCT_DEVELOPMENT_ROADMAP.md")
    return slide

def create_slide_6_financial_projections(prs):
    """Slide 6: Financial Projections."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_headline(slide, "Financial Projections Show Strong ROI by Year 2")

    left_top = MARGIN + Inches(0.8)

    # Left column - Key metrics
    add_section_header(slide, MARGIN, left_top, "REVENUE MODEL ASSUMPTIONS")
    add_body_text(
        slide, MARGIN, left_top + Inches(0.3), LEFT_COL_WIDTH,
        "Conservative scenario based on 15% capture of quality-seeking segment. Hero product pricing $49-89 with 65% gross margin target at scale."
    )

    metrics_data = [
        ("Break-Even Point", "Month 14"),
        ("Year 3 Cumulative Profit", "$2.94M"),
        ("Year 3 Revenue", "$3.15M"),
        ("Year 3 ROI", "1265%")
    ]

    metrics_top = left_top + Inches(1.0)
    for i, (label, value) in enumerate(metrics_data):
        y = metrics_top + i * Inches(0.4)
        add_text_box(
            slide, MARGIN, y, LEFT_COL_WIDTH * 0.6, Inches(0.2),
            label, font_size=10, color=TEXT_LIGHT
        )
        add_text_box(
            slide, MARGIN + LEFT_COL_WIDTH * 0.6, y, LEFT_COL_WIDTH * 0.4, Inches(0.2),
            value, font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT
        )

    # Right column - Projections table
    right_left = MARGIN + LEFT_COL_WIDTH + GAP

    add_text_box(
        slide, right_left, left_top, RIGHT_COL_WIDTH, Inches(0.25),
        "3-YEAR FINANCIAL FORECAST", font_size=12, bold=True, color=CHARCOAL
    )

    # Table headers
    headers = ["Metric", "Year 1", "Year 2", "Year 3"]
    header_y = left_top + Inches(0.35)
    col_width = RIGHT_COL_WIDTH / 4

    for i, header in enumerate(headers):
        add_text_box(
            slide, right_left + i * col_width, header_y, col_width, Inches(0.2),
            header, font_size=9, bold=True, color=CHARCOAL,
            align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        )

    # Table rows
    rows = [
        ("Units", "10,000", "28,000", "45,000"),
        ("Revenue", "$640K", "$1.96M", "$3.15M"),
        ("Gross Profit", "$416K", "$1.27M", "$2.05M"),
        ("Investment", "$900K", "$200K", "$150K"),
        ("ROI", "-46%", "537%", "1265%")
    ]

    row_y = header_y + Inches(0.35)
    for row_data in rows:
        is_roi = row_data[0] == "ROI"
        font_color = ACCENT if is_roi else TEXT

        for i, value in enumerate(row_data):
            add_text_box(
                slide, right_left + i * col_width, row_y, col_width, Inches(0.2),
                value, font_size=10 if not is_roi else 11,
                bold=is_roi, color=font_color,
                align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
            )
        row_y += Inches(0.28)

    add_footer(slide, "Financial model from 01_EXECUTIVE_BRIEFING.md conservative scenario")
    return slide

def create_slide_7_next_steps(prs):
    """Slide 7: Immediate Action Items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_headline(slide, "Immediate Action Items: 12-Week Launch Plan")

    left_top = MARGIN + Inches(0.8)

    # Action items with timeline
    actions = [
        {
            "week": "WEEK 1-2",
            "action": "Technical Validation",
            "tasks": [
                "Validate VHB adhesion on top 10 garage surface types",
                "Test painted drywall, concrete, wood, metal surfaces",
                "ASTM D3330 compliance testing"
            ]
        },
        {
            "week": "WEEK 3-4",
            "action": "Consumer Concept Testing",
            "tasks": [
                "3D printed prototype development (50 units)",
                "30 in-home installations documented",
                "Installation time and ease assessment"
            ]
        },
        {
            "week": "WEEK 5-8",
            "action": "Retailer Partnership",
            "tasks": [
                "Home Depot exclusive launch discussions",
                "50 store test market planning",
                "End-cap display program development"
            ]
        },
        {
            "week": "WEEK 9-12",
            "action": "Pilot Production",
            "tasks": [
                "500 unit pilot production run",
                "Field testing in actual use conditions",
                "Package design finalization"
            ]
        }
    ]

    for i, action_data in enumerate(actions):
        y = left_top + i * Inches(0.9)

        # Week label (left accent)
        week_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN, y, Inches(1.0), Inches(0.7)
        )
        week_box.fill.solid()
        week_box.fill.fore_color.rgb = ACCENT
        week_box.line.fill.background()

        add_text_box(
            slide, MARGIN + Inches(0.1), y + Inches(0.25), Inches(0.8), Inches(0.2),
            action_data["week"], font_size=11, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER
        )

        # Action title
        add_text_box(
            slide, MARGIN + Inches(1.2), y + Inches(0.05), Inches(7.6), Inches(0.25),
            action_data["action"], font_size=13, bold=True, color=CHARCOAL
        )

        # Tasks (bullets)
        tasks_text = "\n".join(f"• {task}" for task in action_data["tasks"])
        add_text_box(
            slide, MARGIN + Inches(1.2), y + Inches(0.32), Inches(7.6), Inches(0.5),
            tasks_text, font_size=10, color=TEXT, line_spacing=1.4
        )

    add_footer(slide, "Action plan from 01_EXECUTIVE_BRIEFING.md | Timeline assumes immediate approval")
    return slide

def main():
    """Generate complete garage organizer deck."""
    print("🚀 Creating 3M Garage Organization Category Intelligence Deck")
    print("=" * 60)

    prs = create_presentation()

    print("📊 Generating slides...")
    create_slide_1_executive_summary(prs)
    print("  ✅ Slide 1: Executive Summary")

    create_slide_2_market_opportunity(prs)
    print("  ✅ Slide 2: Market Opportunity")

    create_slide_3_competitive_positioning(prs)
    print("  ✅ Slide 3: Competitive Positioning")

    create_slide_4_technology_advantage(prs)
    print("  ✅ Slide 4: Technology Advantage")

    create_slide_5_product_roadmap(prs)
    print("  ✅ Slide 5: Product Roadmap")

    create_slide_6_financial_projections(prs)
    print("  ✅ Slide 6: Financial Projections")

    create_slide_7_next_steps(prs)
    print("  ✅ Slide 7: Next Steps")

    # Save presentation
    output_file = "Garage_Organizers_Category_Intelligence_CLIENT_DECK.pptx"
    prs.save(output_file)

    print("=" * 60)
    print(f"✅ DECK CREATED SUCCESSFULLY")
    print(f"📁 File: {output_file}")
    print(f"📊 Slides: {len(prs.slides)}")
    print(f"🎨 Design: offbrain BOLD (Charcoal/Teal)")
    print(f"📐 Format: 16:9 (1920x1080px)")
    print()
    print("💡 Next steps:")
    print("   1. Open in PowerPoint/Keynote")
    print("   2. Review content and formatting")
    print("   3. Install Inter font if needed (or use Helvetica Neue)")
    print("   4. Export to PDF for client distribution")

if __name__ == "__main__":
    main()
