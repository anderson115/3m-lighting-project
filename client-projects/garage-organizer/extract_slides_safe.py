#!/usr/bin/env python3
"""Extract all slide content from PowerPoint presentation"""
from pptx import Presentation
import json
import sys

prs = Presentation('06-final-deliverables/V3-3m_garage_organization_strategy_20251113182641.pptx')
slides = []

for i, slide in enumerate(prs.slides, 1):
    slide_data = {'slide_number': i, 'title': '', 'content': []}

    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            text = shape.text.strip()
            # First text box is usually title
            if not slide_data['title'] and len(text) < 200:
                slide_data['title'] = text
            else:
                slide_data['content'].append(text)

    slides.append(slide_data)

with open('06-final-deliverables/slides_extracted.json', 'w') as f:
    json.dump(slides, f, indent=2)

print(f"Extracted {len(slides)} slides")
for slide in slides:
    print(f"\nSlide {slide['slide_number']}: {slide['title'][:80] if slide['title'] else '(no title)'}")
