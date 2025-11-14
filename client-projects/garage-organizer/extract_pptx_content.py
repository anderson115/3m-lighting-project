#!/usr/bin/env python3
"""
Extract all text content from PowerPoint for audit
"""

from pptx import Presentation
import json

pptx_path = "/Users/anderson115/Downloads/3m_garage_organization_strategy_20251113171531.pptx"
prs = Presentation(pptx_path)

slides_content = []

for idx, slide in enumerate(prs.slides, 1):
    slide_data = {
        'slide_number': idx,
        'title': '',
        'content': [],
        'notes': ''
    }

    # Extract title
    if slide.shapes.title:
        slide_data['title'] = slide.shapes.title.text

    # Extract all text from shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text and text != slide_data['title']:  # Don't duplicate title
                slide_data['content'].append(text)

    # Extract notes
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            slide_data['notes'] = notes_slide.notes_text_frame.text.strip()

    slides_content.append(slide_data)

    # Print slide summary
    print(f"Slide {idx}: {slide_data['title']}")
    if slide_data['content']:
        print(f"  Content items: {len(slide_data['content'])}")

# Save to JSON
output_file = "pptx_extracted_content.json"
with open(output_file, 'w', encoding='utf-8') as f:
    json.dump(slides_content, f, indent=2, ensure_ascii=False)

print(f"\n✅ Extracted content from {len(slides_content)} slides")
print(f"✅ Saved to: {output_file}")
